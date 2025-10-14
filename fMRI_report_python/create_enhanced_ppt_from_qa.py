#!/usr/bin/env python
"""
Enhanced QA PowerPoint generator with comprehensive metrics extraction
Usage: python create_enhanced_ppt_from_qa.py --qa_parent_dir /path/to/parent/ --func_dir /path/to/func/
"""

import argparse
import os
import json
from datetime import datetime
import numpy as np

# PowerPoint libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed.")

# Neuroimaging libraries
try:
    import nibabel as nib
    NIBABEL_AVAILABLE = True
except ImportError:
    NIBABEL_AVAILABLE = False
    print("⚠️  nibabel not available for detailed metrics extraction")

def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception as e:
        print(f"Warning: Could not read TR from JSON: {e}")
    return None

def calculate_comprehensive_metrics(nifti_file, qa_dir):
    """Calculate comprehensive metrics from NIFTI file and QA outputs"""
    metrics = {}
    
    if not NIBABEL_AVAILABLE:
        return metrics
    
    try:
        # Load original data
        img = nib.load(nifti_file)
        data = img.get_fdata()
        
        metrics['data_shape'] = data.shape
        metrics['voxel_size'] = img.header.get_zooms()
        
        # Basic statistics
        mean_img = np.mean(data, axis=-1)
        metrics['mean_signal'] = float(np.mean(mean_img))
        metrics['std_signal'] = float(np.std(mean_img))
        
        # Get slice index (middle slice)
        slice_index = data.shape[2] // 2
        metrics['slice_index'] = slice_index
        
        # Calculate brain mask (simple threshold)
        brain_threshold = 0.05 * np.max(mean_img)
        brain_mask = mean_img > brain_threshold
        
        if np.any(brain_mask):
            valid_voxels = mean_img[brain_mask]
            metrics['mean_volume_std'] = float(np.std(valid_voxels))
            
        # Load QA-derived metrics from NIFTI files
        isnr_file = os.path.join(qa_dir, 'isnr.nii.gz')
        tsnr_file = os.path.join(qa_dir, 'tsnr.nii.gz')
        
        if os.path.exists(isnr_file):
            isnr_img = nib.load(isnr_file)
            isnr_data = isnr_img.get_fdata()
            if isnr_data.ndim > 3:
                isnr_data = isnr_data[:, :, :, 0]  # Take first volume if 4D
            metrics['isnr'] = float(np.mean(isnr_data[isnr_data > 0]))
            
        if os.path.exists(tsnr_file):
            tsnr_img = nib.load(tsnr_file)
            tsnr_data = tsnr_img.get_fdata()
            metrics['tsnr'] = float(np.mean(tsnr_data[tsnr_data > 0]))
            
            # If we have brain mask, calculate masked tSNR
            if np.any(brain_mask) and tsnr_data.shape[:3] == brain_mask.shape:
                masked_tsnr = tsnr_data[brain_mask]
                metrics['masked_tsnr'] = float(np.mean(masked_tsnr[masked_tsnr > 0]))
        
        # Noise estimation (using last volume if available)
        if data.shape[-1] > 1:
            noise_volume = data[:, :, :, -1]
            noise_volume = np.nan_to_num(noise_volume, nan=0.0)
            masked_noise = noise_volume * brain_mask
            if np.any(masked_noise):
                metrics['noise_value'] = float(np.mean(masked_noise[masked_noise > 0]))
        
    except Exception as e:
        print(f"Error calculating metrics for {nifti_file}: {e}")
    
    return metrics

def extract_enhanced_metrics(qa_output_dirs, func_dir=None):
    """Extract comprehensive metrics from QA directories and original data"""
    metrics_list = []
    
    for qa_dir in qa_output_dirs:
        if not os.path.exists(qa_dir):
            continue
            
        # Extract filename from directory name
        dirname = os.path.basename(qa_dir)
        if dirname.startswith('qa_output_'):
            filename = dirname[10:]  # Remove 'qa_output_' prefix
        else:
            filename = dirname
            
        # Initialize metrics dictionary
        metrics = {
            'filename': filename,
            'qa_dir': qa_dir,
            'tr': 'Unknown',
            'ernst_factor': None
        }
        
        # Try to find corresponding NIFTI file and extract comprehensive metrics
        if func_dir:
            for ext in ['.nii.gz', '.nii']:
                nifti_file = os.path.join(func_dir, filename + ext)
                if os.path.exists(nifti_file):
                    # Get TR and Ernst factor
                    tr_value = get_tr_from_json(nifti_file)
                    if tr_value:
                        metrics['tr'] = tr_value
                        # Calculate Ernst scaling factor
                        if tr_value <= 0.7:
                            metrics['ernst_factor'] = 0.5745
                        elif tr_value <= 1.0:
                            metrics['ernst_factor'] = 0.7071
                        elif tr_value <= 1.5:
                            metrics['ernst_factor'] = 0.8155
                        else:
                            metrics['ernst_factor'] = 1.0
                        
                        # Calculate tSNR per unit time
                        if 'tsnr' in metrics and metrics['tsnr']:
                            metrics['tsnr_per_unit_time'] = metrics['tsnr'] / np.sqrt(tr_value)
                    
                    # Get comprehensive metrics
                    comprehensive_metrics = calculate_comprehensive_metrics(nifti_file, qa_dir)
                    metrics.update(comprehensive_metrics)
                    
                    # Calculate tSNR per unit time if we have both values
                    if tr_value and 'tsnr' in metrics and metrics['tsnr']:
                        metrics['tsnr_per_unit_time'] = metrics['tsnr'] / np.sqrt(tr_value)
                    
                    break
        
        metrics_list.append(metrics)
    
    return metrics_list

def create_enhanced_qa_powerpoint(qa_output_dirs, presentation_path, summary_data):
    """Create enhanced PowerPoint with comprehensive metrics"""
    if not PPTX_AVAILABLE:
        print("❌ Cannot create PowerPoint - python-pptx not installed")
        return False
    
    print("🎨 Creating enhanced PowerPoint presentation...")
    
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Enhanced fMRI Quality Assessment Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{len(qa_output_dirs)} datasets with comprehensive metrics"
    
    # Enhanced summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Comprehensive QA Metrics Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Dataset Overview with Key QA Metrics:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for i, data in enumerate(summary_data):
        p = tf.add_paragraph()
        p.text = f"\n{i+1}. {data['filename'][:60]}{'...' if len(data['filename']) > 60 else ''}"
        p.font.size = Pt(11)
        p.font.bold = True
        
        # TR and Ernst info
        tr_text = f"   • TR: {data['tr']}s"
        if data.get('ernst_factor'):
            tr_text += f" | Ernst scaling: {data['ernst_factor']:.4f}"
        p = tf.add_paragraph()
        p.text = tr_text
        p.font.size = Pt(9)
        
        # SNR metrics
        snr_text = "   • "
        if data.get('isnr'):
            snr_text += f"iSNR: {data['isnr']:.2f}"
        if data.get('tsnr'):
            snr_text += f" | tSNR: {data['tsnr']:.2f}"
        if data.get('tsnr_per_unit_time'):
            snr_text += f" | tSNR/√time: {data['tsnr_per_unit_time']:.2f}"
        
        if snr_text != "   • ":
            p = tf.add_paragraph()
            p.text = snr_text
            p.font.size = Pt(9)
        
        # Data characteristics
        char_text = "   • "
        if data.get('data_shape'):
            shape = data['data_shape']
            if len(shape) >= 4:
                char_text += f"Shape: {shape[0]}×{shape[1]}×{shape[2]}×{shape[3]}"
        if data.get('voxel_size'):
            vox_size = data['voxel_size']
            char_text += f" | Voxel: {vox_size[0]:.1f}×{vox_size[1]:.1f}×{vox_size[2]:.1f}mm"
        
        if char_text != "   • ":
            p = tf.add_paragraph()
            p.text = char_text
            p.font.size = Pt(9)
    
    # Detailed slides for each dataset
    for i, (qa_dir, data) in enumerate(zip(qa_output_dirs, summary_data)):
        if not os.path.exists(qa_dir):
            continue
            
        # Enhanced dataset overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Dataset {i+1}: Detailed Metrics"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = f"{data['filename'][:70]}{'...' if len(data['filename']) > 70 else ''}"
        p.font.size = Pt(12)
        p.font.bold = True
        
        # Build comprehensive metrics text
        metrics_text = "\nAcquisition Parameters:\n"
        if data['tr'] != 'Unknown':
            metrics_text += f"• TR (Repetition Time): {data['tr']}s\n"
        if data.get('ernst_factor'):
            metrics_text += f"• Ernst scaling factor: {data['ernst_factor']:.4f} for TR = {data['tr']}s\n"
        
        if data.get('data_shape'):
            shape = data['data_shape']
            if len(shape) >= 4:
                metrics_text += f"• Image data shape: {shape[0]} × {shape[1]} × {shape[2]} × {shape[3]} voxels\n"
        
        if data.get('voxel_size'):
            vox = data['voxel_size']
            metrics_text += f"• Voxel size: {vox[0]:.2f} × {vox[1]:.2f} × {vox[2]:.2f} mm\n"
        
        if data.get('slice_index'):
            metrics_text += f"• Central slice index: {data['slice_index']}\n"
        
        metrics_text += "\nQuality Metrics:\n"
        if data.get('isnr'):
            metrics_text += f"• Image SNR (iSNR): {data['isnr']:.2f}\n"
        if data.get('noise_value'):
            metrics_text += f"• Noise value used for iSNR: {data['noise_value']:.2f}\n"
        if data.get('tsnr'):
            metrics_text += f"• Temporal SNR (tSNR): {data['tsnr']:.2f}\n"
        if data.get('masked_tsnr'):
            metrics_text += f"• Masked tSNR: {data['masked_tsnr']:.2f}\n"
        if data.get('tsnr_per_unit_time'):
            metrics_text += f"• tSNR per unit time: {data['tsnr_per_unit_time']:.2f}\n"
        if data.get('mean_volume_std'):
            metrics_text += f"• Mean volume std: {data['mean_volume_std']:.2f}\n"
        
        # Count available images
        image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')] if os.path.exists(qa_dir) else []
        metrics_text += f"\nQA Output:\n"
        metrics_text += f"• Generated {len(image_files)} QA visualization images\n"
        metrics_text += f"• Output directory: {os.path.basename(qa_dir)}"
        
        p = tf.add_paragraph()
        p.text = metrics_text
        p.font.size = Pt(10)
        
        # Add image slides (reuse existing logic)
        if os.path.exists(qa_dir):
            image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
            
            image_groups = {
                'Mean Images': [f for f in image_files if 'mean' in f.lower()],
                'SNR Analysis': [f for f in image_files if any(x in f.lower() for x in ['isnr', 'tsnr', 'snr'])],
                'Montage Views': [f for f in image_files if 'montage' in f.lower()],
                'Noise Analysis': [f for f in image_files if 'noise' in f.lower()],
                'Time Series': [f for f in image_files if any(x in f.lower() for x in ['ts_', 'timeseries'])],
                'Other': [f for f in image_files if not any(keyword in f.lower() for keyword in ['mean', 'snr', 'montage', 'noise', 'ts_', 'timeseries'])]
            }
            
            for category, images in image_groups.items():
                if not images:
                    continue
                    
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"Dataset {i+1}: {category}"
                title_para.font.size = Pt(18)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                y_pos = Inches(1.0)
                for j, img_file in enumerate(images[:4]):
                    img_path = os.path.join(qa_dir, img_file)
                    if os.path.exists(img_path):
                        try:
                            left = Inches(0.5 + (j % 2) * 4.75)
                            if j >= 2:
                                y_pos = Inches(4.0)
                            
                            pic = slide.shapes.add_picture(img_path, left, y_pos)
                            
                            # Scale to fit
                            max_width = Inches(4.5)
                            max_height = Inches(3)
                            
                            if pic.width > max_width:
                                ratio = max_width / pic.width
                                pic.width = max_width
                                pic.height = int(pic.height * ratio)
                            
                            if pic.height > max_height:
                                ratio = max_height / pic.height
                                pic.height = max_height
                                pic.width = int(pic.width * ratio)
                            
                            caption_shape = slide.shapes.add_textbox(left, y_pos + pic.height, pic.width, Inches(0.5))
                            caption_frame = caption_shape.text_frame
                            caption_para = caption_frame.paragraphs[0]
                            caption_para.text = img_file.replace('.png', '').replace('_', ' ').title()
                            caption_para.font.size = Pt(9)
                            caption_para.alignment = PP_ALIGN.CENTER
                            
                        except Exception as e:
                            print(f"⚠️  Could not add image {img_file}: {e}")
    
    # Save presentation
    try:
        prs.save(presentation_path)
        print(f"✅ Enhanced PowerPoint saved: {presentation_path}")
        return True
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(description='Create enhanced PowerPoint from QA directories with comprehensive metrics')
    parser.add_argument('--qa_parent_dir', required=True, help='Parent directory containing qa_output_* folders')
    parser.add_argument('--func_dir', required=True, help='Func directory with original NIFTI files')
    parser.add_argument('--output_name', default='Enhanced_QA_Report.pptx', help='PowerPoint filename')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.qa_parent_dir):
        print(f"❌ QA parent directory does not exist: {args.qa_parent_dir}")
        return 1
        
    if not os.path.exists(args.func_dir):
        print(f"❌ Func directory does not exist: {args.func_dir}")
        return 1
    
    # Find QA output directories
    qa_dirs = []
    for item in os.listdir(args.qa_parent_dir):
        if item.startswith('qa_output_') and os.path.isdir(os.path.join(args.qa_parent_dir, item)):
            qa_dirs.append(os.path.join(args.qa_parent_dir, item))
    
    if not qa_dirs:
        print(f"❌ No qa_output_* directories found in {args.qa_parent_dir}")
        return 1
    
    print(f"📁 Found {len(qa_dirs)} QA output directories")
    
    # Extract comprehensive metrics
    summary_data = extract_enhanced_metrics(qa_dirs, args.func_dir)
    
    # Create enhanced PowerPoint
    ppt_path = os.path.join(args.qa_parent_dir, args.output_name)
    success = create_enhanced_qa_powerpoint(qa_dirs, ppt_path, summary_data)
    
    if success:
        print(f"🎉 Enhanced PowerPoint report created successfully!")
        print(f"📄 File: {ppt_path}")
    else:
        print("❌ Failed to create enhanced PowerPoint report")
    
    return 0 if success else 1

if __name__ == "__main__":
    exit(main())