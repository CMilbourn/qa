#!/usr/bin/env python
"""
QA Analysis with PowerPoint Report Generation (V2)
Automatically creates a PowerPoint presentation with all QA images and metrics
Usage: python run_qa_with_ppt_V2.py --func_dir /path/to/func/ [options]
"""

import argparse
import os
import sys
from pathlib import Path
import time
from datetime import datetime
import json

# Add the current directory to Python path for imports
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

# Import the main QA functions
from qa_with_metrics import run_qa_with_metrics
import numpy as np
from glob import glob

# PowerPoint libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed. Install with: pip install python-pptx")

def create_qa_powerpoint(output_directories, presentation_path, summary_data):
    """
    Create a PowerPoint presentation with QA results
    
    Parameters:
    -----------
    output_directories : list
        List of QA output directories
    presentation_path : str
        Path where to save the PowerPoint file
    summary_data : list
        List of dictionaries with QA metrics for each file
    """
    if not PPTX_AVAILABLE:
        print("❌ Cannot create PowerPoint - python-pptx not installed")
        return False
    
    print("🎨 Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assessment Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{len(output_directories)} datasets analyzed"
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "QA Summary"
    
    # Add summary table
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Header
    p = tf.paragraphs[0]
    p.text = "Dataset Summary:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # Add metrics for each dataset
    for i, data in enumerate(summary_data):
        p = tf.add_paragraph()
        p.text = f"\n{i+1}. {data['filename'][:50]}..."
        p.font.size = Pt(12)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"   • TR: {data['tr']}s | iSNR: {data['isnr']:.2f} | tSNR: {data['tsnr']:.2f}"
        p.font.size = Pt(10)
        
        p = tf.add_paragraph()
        p.text = f"   • Data shape: {data['shape']} | Ernst scaling: {data['ernst']}"
        p.font.size = Pt(10)
    
    # Create slides for each dataset
    for i, (output_dir, data) in enumerate(zip(output_directories, summary_data)):
        
        # Dataset overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Dataset {i+1}: QA Overview"
        
        # Add dataset info
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = f"File: {data['filename']}"
        p.font.size = Pt(14)
        p.font.bold = True
        
        metrics_text = f"""
Key Metrics:
• TR (Repetition Time): {data['tr']}s
• Ernst Scaling Factor: {data['ernst']}
• Image SNR (iSNR): {data['isnr']:.2f}
• Temporal SNR (tSNR): {data['tsnr']:.2f}
• tSNR per unit time: {data['tsnr_unit']:.2f}
• Data dimensions: {data['shape']}
• Mean SSN: {data['ssn']:.1f}

Output Directory:
{os.path.basename(output_dir)}
        """
        
        p = tf.add_paragraph()
        p.text = metrics_text
        p.font.size = Pt(12)
        
        # Image slides for this dataset
        image_files = {
            'Mean Images': ['Mean_image.png', 'mean_montage.png'],
            'SNR Analysis': ['iSNR_sag.png', 'iSNR_cor.png', 'tSNR_sag.png', 'tSNR_cor.png'],
            'Montage Views': ['isnr_montage.png', 'tSNR_montage.png'],
            'Noise Analysis': ['masked_noise.png', 'noise_volume_montage.png'],
            'Time Series': ['TS_images.png', 'tSNR_w_ROI_images.png'],
            'tSNR Comparison': ['tSNR_raw.png', 'tSNR_per_unit_time.png'],
            'Spatial Noise': ['SSN.png']
        }
        
        for category, images in image_files.items():
            # Create slide for this category
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"Dataset {i+1}: {category}"
            title_para.font.size = Pt(20)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add images
            y_pos = Inches(1.0)
            images_added = 0
            
            for img_file in images:
                img_path = os.path.join(output_dir, img_file)
                if os.path.exists(img_path):
                    try:
                        # Determine image position
                        if len(images) == 1:
                            # Single image - center it
                            left = Inches(2)
                            width = Inches(6)
                            height = Inches(4.5)
                        elif len(images) == 2:
                            # Two images side by side
                            left = Inches(0.5 + (images_added * 4.5))
                            width = Inches(4)
                            height = Inches(3)
                        else:
                            # Multiple images - smaller grid
                            left = Inches(0.5 + (images_added % 2) * 4.5)
                            if images_added >= 2:
                                y_pos = Inches(4.0)
                            width = Inches(4)
                            height = Inches(2.5)
                        
                        # Add image
                        slide.shapes.add_picture(img_path, left, y_pos, width, height)
                        images_added += 1
                        
                        # Add image caption
                        caption_shape = slide.shapes.add_textbox(left, y_pos + height, width, Inches(0.3))
                        caption_frame = caption_shape.text_frame
                        caption_para = caption_frame.paragraphs[0]
                        caption_para.text = img_file.replace('.png', '').replace('_', ' ').title()
                        caption_para.font.size = Pt(10)
                        caption_para.alignment = PP_ALIGN.CENTER
                        
                    except Exception as e:
                        print(f"⚠️  Could not add image {img_file}: {e}")
    
    # Save presentation
    try:
        prs.save(presentation_path)
        print(f"✅ PowerPoint saved: {presentation_path}")
        return True
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def run_qa_with_powerpoint(func_dir, pattern='*fmri*', extension='.nii.gz', 
                          output_name=None, include_ppt=True):
    """
    Run QA analysis and create PowerPoint presentation
    
    Parameters:
    -----------
    func_dir : str
        Path to the func directory containing the BOLD files
    pattern : str
        Filename pattern to match
    extension : str
        File extension
    output_name : str
        Custom name for the PowerPoint file
    include_ppt : bool
        Whether to create PowerPoint presentation
    """
    
    print(f"🔍 fMRI QA Analysis with PowerPoint Generation")
    print(f"📁 Directory: {func_dir}")
    print(f"🔎 Pattern: {pattern}{extension}")
    
    # Run QA analysis with metrics collection
    results = run_qa_with_metrics(func_dir, pattern, extension)
    
    if not results:
        print(f"❌ No successful analyses to create PowerPoint from")
        return False
    
    # Extract data for PowerPoint
    output_directories = []
    summary_data = []
    
    for output_dir, metrics in results:
        output_directories.append(output_dir)
        
        # Format metrics for PowerPoint
        summary_info = {
            'filename': metrics.get('filename', 'Unknown'),
            'tr': metrics.get('tr', 'Unknown'),
            'ernst': metrics.get('ernst_factor', 1.0),
            'shape': f"{metrics['shape'][0]}×{metrics['shape'][1]}×{metrics['shape'][2]}×{metrics['shape'][3]}" if 'shape' in metrics else 'Unknown',
            'isnr': metrics.get('isnr', 0.0),
            'tsnr': metrics.get('tsnr', 0.0),
            'tsnr_unit': metrics.get('tsnr_per_unit_time', 0.0),
            'ssn': metrics.get('ssn', 0.0),
            'processing_time': 0  # Not tracked in new version
        }
        
        summary_data.append(summary_info)
    
    success_count = len(results)
    print(f"\n{'='*80}")
    print(f"🎉 QA Analysis Complete: {success_count} files processed")
    
    # Create PowerPoint presentation
    if include_ppt and success_count > 0 and PPTX_AVAILABLE:
        if output_name is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            output_name = f"fMRI_QA_Report_{timestamp}.pptx"
        
        # Save in parent directory of func folder
        parent_dir = os.path.dirname(func_dir)
        ppt_path = os.path.join(parent_dir, output_name)
        
        print(f"\n🎨 Creating PowerPoint presentation...")
        ppt_success = create_qa_powerpoint(output_directories, ppt_path, summary_data)
        
        if ppt_success:
            print(f"📊 PowerPoint report created: {ppt_path}")
        else:
            print(f"❌ Failed to create PowerPoint report")
    elif not PPTX_AVAILABLE:
        print(f"\n⚠️  PowerPoint generation skipped - install python-pptx to enable")
    
    print(f"{'='*80}")
    return success_count > 0

def main():
    parser = argparse.ArgumentParser(
        description='Run fMRI QA analysis with PowerPoint report generation',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python run_qa_with_ppt_V2.py --func_dir /path/to/func/
  python run_qa_with_ppt_V2.py --func_dir /path/to/func/ --output_name "MyQA_Report.pptx"
  python run_qa_with_ppt_V2.py --func_dir /path/to/func/ --pattern "*bold*" --no_ppt
        """
    )
    
    parser.add_argument('--func_dir', required=True,
                        help='Path to func directory with BOLD files')
    parser.add_argument('--pattern', default='*fmri*',
                        help='Filename pattern (default: *fmri*)')
    parser.add_argument('--extension', default='.nii.gz',
                        help='File extension (default: .nii.gz)')
    parser.add_argument('--output_name', 
                        help='Custom PowerPoint filename')
    parser.add_argument('--no_ppt', action='store_true',
                        help='Skip PowerPoint generation')
    
    args = parser.parse_args()
    
    success = run_qa_with_powerpoint(
        func_dir=args.func_dir,
        pattern=args.pattern, 
        extension=args.extension,
        output_name=args.output_name,
        include_ppt=not args.no_ppt
    )
    
    return 0 if success else 1

if __name__ == "__main__":
    # Check for dependencies
    if not PPTX_AVAILABLE:
        print("Installing python-pptx for PowerPoint generation...")
        try:
            import subprocess
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
            print("✅ python-pptx installed successfully!")
            # Restart to import the module
            os.execv(sys.executable, ['python'] + sys.argv)
        except Exception as e:
            print(f"❌ Failed to install python-pptx: {e}")
            print("Please install manually: pip install python-pptx")
            sys.exit(1)
    
    sys.exit(main())