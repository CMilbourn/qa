#!/usr/bin/env python3
"""
Create PowerPoint from NIFTI files (_nodummies) with matplotlib visualization
Renders 3-view (axial/sagittal/coronal) slices of NIFTI data into PowerPoint
"""

import argparse
import os
from pathlib import Path
import json
import tempfile
import shutil

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed.")

try:
    import nibabel as nib
    import matplotlib.pyplot as plt
    import numpy as np
    NIBABEL_AVAILABLE = True
except ImportError:
    NIBABEL_AVAILABLE = False
    print("⚠️  nibabel or matplotlib not installed.")


def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = str(nifti_path).replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception:
        pass
    return None


def create_3view_visualization(nifti_path, output_png_path):
    """
    Create a 3-view visualization (axial, sagittal, coronal) using matplotlib
    Shows a middle slice of the volume
    """
    if not NIBABEL_AVAILABLE:
        print("⚠️  nibabel required for visualization")
        return False
    
    try:
        img = nib.load(nifti_path)
        data = img.get_fdata()
        
        # Handle 4D data (fMRI) - use middle timepoint
        if len(data.shape) == 4:
            vol_idx = data.shape[3] // 2
            vol_data = data[:, :, :, vol_idx]
        else:
            vol_data = data
        
        # Get middle slices
        x_mid = vol_data.shape[0] // 2
        y_mid = vol_data.shape[1] // 2
        z_mid = vol_data.shape[2] // 2
        
        fig, axes = plt.subplots(1, 3, figsize=(15, 5), dpi=100)
        fig.patch.set_facecolor('white')
        
        # Axial (z-slice)
        axial_slice = vol_data[:, :, z_mid]
        im1 = axes[0].imshow(axial_slice.T, cmap='gray', origin='lower')
        axes[0].set_title('Axial', fontsize=12, fontweight='bold')
        axes[0].axis('off')
        
        # Sagittal (x-slice)
        sagittal_slice = vol_data[x_mid, :, :]
        im2 = axes[1].imshow(sagittal_slice.T, cmap='gray', origin='lower')
        axes[1].set_title('Sagittal', fontsize=12, fontweight='bold')
        axes[1].axis('off')
        
        # Coronal (y-slice)
        coronal_slice = vol_data[:, y_mid, :]
        im3 = axes[2].imshow(coronal_slice.T, cmap='gray', origin='lower')
        axes[2].set_title('Coronal', fontsize=12, fontweight='bold')
        axes[2].axis('off')
        
        plt.tight_layout()
        plt.savefig(output_png_path, dpi=100, bbox_inches='tight', facecolor='white')
        plt.close(fig)
        
        return True
        
    except Exception as e:
        print(f"⚠️  Error creating visualization: {e}")
        return False


def find_nifti_files(base_dir):
    """
    Find all NIFTI files with _nodummies in QA subdirectories
    Returns list of (qa_dir, nifti_path) tuples
    """
    nifti_files = []
    base_path = Path(base_dir)
    
    # Find all subdirectories (QA output directories)
    for qa_dir in sorted(base_path.iterdir()):
        if qa_dir.is_dir():
            # Look for *_nodummies.nii.gz files
            for nifti_file in sorted(qa_dir.glob('*_nodummies.nii.gz')):
                nifti_files.append((qa_dir.name, str(nifti_file)))
            
            # Also look for other .nii.gz files if no nodummies found
            if not any(qa_dir.name in item[0] for item in nifti_files):
                for nifti_file in sorted(qa_dir.glob('*.nii.gz')):
                    if '_nodummies' not in nifti_file.name:
                        nifti_files.append((qa_dir.name, str(nifti_file)))
                        break  # Only take first file per directory
    
    return nifti_files


def add_metrics_table(slide, nifti_files_with_metrics):
    """Add a summary table with metrics for all files"""
    if not nifti_files_with_metrics:
        return
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = "Scan Metrics Summary"
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    
    # Create table
    rows = len(nifti_files_with_metrics) + 1
    cols = 4
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9)
    height = Inches(0.3 * rows)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ['Scan Name', 'TR (s)', 'Shape', 'File']
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
    
    # Data rows
    for row_idx, (qa_dir_name, nifti_path, tr, shape) in enumerate(nifti_files_with_metrics, 1):
        # Scan name
        cell = table_shape.cell(row_idx, 0)
        cell.text = qa_dir_name[:40]  # Truncate long names
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        
        # TR
        cell = table_shape.cell(row_idx, 1)
        cell.text = f"{tr:.3f}" if tr else "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Shape
        cell = table_shape.cell(row_idx, 2)
        if shape:
            shape_str = '×'.join(str(s) for s in shape[:3])  # Show first 3 dims
            cell.text = shape_str
        else:
            cell.text = "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # File name
        cell = table_shape.cell(row_idx, 3)
        fname = Path(nifti_path).name[:30]
        cell.text = fname
        cell.text_frame.paragraphs[0].font.size = Pt(7)


def create_powerpoint_from_nifti(input_dir, output_pptx):
    """
    Create PowerPoint presentation from NIFTI files
    """
    if not PPTX_AVAILABLE or not NIBABEL_AVAILABLE:
        print("❌ Required packages not available")
        return False
    
    print(f"📁 Input directory: {input_dir}")
    print(f"📄 Output file: {output_pptx}")
    
    # Find all NIFTI files
    nifti_files = find_nifti_files(input_dir)
    
    if not nifti_files:
        print("❌ No NIFTI files found")
        return False
    
    print(f"📊 Found {len(nifti_files)} NIFTI files")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Temporary directory for images
    temp_dir = tempfile.mkdtemp()
    
    try:
        # Create title slide
        title_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(title_slide_layout)
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "NIFTI Data Visualization Report\n(Sub009 - nodummies)"
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Collect metrics for summary
        metrics_data = []
        processed_count = 0
        
        # Process each NIFTI file
        for qa_dir_name, nifti_path in nifti_files:
            try:
                print(f"\n📁 Processing: {qa_dir_name}")
                print(f"   📊 Found NIFTI: {Path(nifti_path).name}")
                
                # Extract metrics
                tr = get_tr_from_json(nifti_path)
                try:
                    img = nib.load(nifti_path)
                    shape = img.shape
                except:
                    shape = None
                
                metrics_data.append((qa_dir_name, nifti_path, tr if tr else 0.0, shape))
                
                # Create visualization
                temp_png = os.path.join(temp_dir, f"{processed_count:03d}_tmp.png")
                print(f"   🎨 Creating 3-view visualization...")
                
                if create_3view_visualization(nifti_path, temp_png):
                    # Add slide with image
                    blank_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_layout)
                    
                    # Add title with scan name
                    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
                    title_frame = title_box.text_frame
                    title_frame.text = f"Scan: {qa_dir_name}"
                    title_frame.paragraphs[0].font.size = Pt(14)
                    title_frame.paragraphs[0].font.bold = True
                    
                    # Add metrics box
                    info_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.3))
                    info_frame = info_box.text_frame
                    info_text = f"TR: {tr:.3f}s" if tr else "TR: N/A"
                    if shape:
                        info_text += f" | Shape: {' × '.join(str(s) for s in shape)}"
                    info_frame.text = info_text
                    info_frame.paragraphs[0].font.size = Pt(10)
                    
                    # Add image
                    img_path = temp_png
                    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9))
                    
                    processed_count += 1
                    print(f"   ✅ Slide added")
                else:
                    print(f"   ❌ Failed to create visualization")
                    
            except Exception as e:
                print(f"   ❌ Error processing {nifti_path}: {e}")
        
        # Add metrics summary slide
        if metrics_data:
            print(f"\n📊 Adding metrics summary slide...")
            summary_slide = prs.slides.add_slide(prs.slide_layouts[6])
            add_metrics_table(summary_slide, metrics_data)
        
        if processed_count == 0:
            print("❌ No visualizations created")
            return False
        
        # Save presentation
        prs.save(output_pptx)
        print(f"\n✅ PowerPoint saved: {output_pptx}")
        print(f"📊 Total slides: {len(prs.slides)} (1 title + {processed_count} images + 1 metrics)")
        
        return True
        
    finally:
        # Clean up temporary directory
        shutil.rmtree(temp_dir, ignore_errors=True)


def main():
    parser = argparse.ArgumentParser(
        description='Create PowerPoint from NIFTI files with 3-view visualization'
    )
    parser.add_argument('input_dir', help='Input directory containing QA subdirectories')
    parser.add_argument('--output', '-o', default='sub009_NIFTI_nodummies.pptx',
                       help='Output PowerPoint file')
    
    args = parser.parse_args()
    
    if not os.path.isdir(args.input_dir):
        print(f"❌ Input directory not found: {args.input_dir}")
        return 1
    
    success = create_powerpoint_from_nifti(args.input_dir, args.output)
    return 0 if success else 1


if __name__ == '__main__':
    exit(main())
