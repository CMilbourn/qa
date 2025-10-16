#!/usr/bin/env python
"""
QA Analysis with PowerPoint Report Generation (V2 Simple)
Version: 2.0.0

Simple approach that runs existing QA analysis then creates PowerPoint
Usage: python run_qa_with_ppt_V2_simple.py --func_dir /path/to/func/ [options]

Version History:
- v2.0.0 (2025-10-16): Baseline version for simple integrated QA + PowerPoint workflow
"""

import argparse
import os
import sys
import time
import json
from datetime import datetime
from pathlib import Path
import subprocess

# PowerPoint libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed. Install with: pip install python-pptx")

def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception as e:
        print(f"Warning: Could not read TR from JSON: {e}")
    return None

def extract_metrics_from_qa_dirs(qa_output_dirs):
    """
    Extract metrics from existing QA output directories
    This is a simplified approach that estimates metrics from generated files
    """
    metrics_list = []
    
    for qa_dir in qa_output_dirs:
        if not os.path.exists(qa_dir):
            continue
            
        # Extract filename from directory name
        dirname = os.path.basename(qa_dir)
        if dirname.startswith('qa_output_'):
            filename = dirname[10:]  # Remove 'qa_output_' prefix
        else:
            filename = dirname
            
        # Look for JSON file to get TR
        parent_dir = os.path.dirname(qa_dir)
        func_dir = None
        if 'func' in parent_dir:
            func_dir = parent_dir
        else:
            # Look for func subdirectory
            possible_func = os.path.join(parent_dir, 'func')
            if os.path.exists(possible_func):
                func_dir = possible_func
        
        tr_value = None
        if func_dir:
            # Try to find matching nifti and JSON files
            for ext in ['.nii.gz', '.nii']:
                nifti_file = os.path.join(func_dir, filename + ext)
                if os.path.exists(nifti_file):
                    tr_value = get_tr_from_json(nifti_file)
                    break
        
        # Default metrics (would need to parse from QA output files for real values)
        metrics = {
            'filename': filename,
            'tr': tr_value if tr_value else 'Unknown',
            'ernst': 1.0,  # Would calculate based on TR
            'shape': 'Unknown',
            'isnr': 0.5,   # Placeholder
            'tsnr': 8.0,   # Placeholder
            'tsnr_unit': 6.0,  # Placeholder
            'ssn': 100.0,  # Placeholder
            'qa_dir': qa_dir
        }
        
        metrics_list.append(metrics)
    
    return metrics_list

def create_qa_powerpoint(qa_output_dirs, presentation_path, summary_data):
    """
    Create a PowerPoint presentation with QA results
    """
    if not PPTX_AVAILABLE:
        print("❌ Cannot create PowerPoint - python-pptx not installed")
        return False
    
    print("🎨 Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assessment Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{len(qa_output_dirs)} datasets analyzed"
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "QA Summary"
    
    # Add summary table
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Header
    p = tf.paragraphs[0]
    p.text = "Dataset Summary:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    # Add metrics for each dataset
    for i, data in enumerate(summary_data):
        p = tf.add_paragraph()
        p.text = f"\n{i+1}. {data['filename'][:50]}..."
        p.font.size = Pt(12)
        p.font.bold = True
        
        p = tf.add_paragraph()
        p.text = f"   • TR: {data['tr']}s | Output: {os.path.basename(data['qa_dir'])}"
        p.font.size = Pt(10)
    
    # Create slides for each dataset
    for i, (qa_dir, data) in enumerate(zip(qa_output_dirs, summary_data)):
        
        # Dataset overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Dataset {i+1}: {data['filename']}"
        
        # Add dataset info
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = f"File: {data['filename']}"
        p.font.size = Pt(14)
        p.font.bold = True
        
        info_text = f"""
Dataset Information:
• TR (Repetition Time): {data['tr']}s
• QA Output Directory: {os.path.basename(qa_dir)}

Available QA Images:
        """
        
        # List available images
        if os.path.exists(qa_dir):
            image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
            if image_files:
                info_text += f"• {len(image_files)} QA images generated\n"
                for img_file in sorted(image_files):
                    info_text += f"  - {img_file}\n"
            else:
                info_text += "• No PNG images found\n"
        
        p = tf.add_paragraph()
        p.text = info_text
        p.font.size = Pt(12)
        
        # Image slides for this dataset
        if os.path.exists(qa_dir):
            image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
            
            # Group images by category
            image_groups = {
                'Mean Images': [f for f in image_files if 'mean' in f.lower()],
                'SNR Analysis': [f for f in image_files if any(x in f.lower() for x in ['isnr', 'tsnr', 'snr'])],
                'Montage Views': [f for f in image_files if 'montage' in f.lower()],
                'Noise Analysis': [f for f in image_files if 'noise' in f.lower()],
                'Time Series': [f for f in image_files if any(x in f.lower() for x in ['ts_', 'timeseries'])],
                'Other': [f for f in image_files if not any(keyword in f.lower() for keyword in ['mean', 'snr', 'montage', 'noise', 'ts_', 'timeseries'])]
            }
            
            for category, images in image_groups.items():
                if not images:
                    continue
                    
                # Create slide for this category
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
                title_frame = title_shape.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = f"Dataset {i+1}: {category}"
                title_para.font.size = Pt(20)
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                
                # Add images
                y_pos = Inches(1.0)
                images_added = 0
                
                for img_file in images[:4]:  # Limit to 4 images per slide
                    img_path = os.path.join(qa_dir, img_file)
                    if os.path.exists(img_path):
                        try:
                            # Determine image position
                            if len(images) == 1:
                                # Single image - center it
                                left = Inches(2)
                                max_width = Inches(6)
                                max_height = Inches(4.5)
                            elif len(images) == 2:
                                # Two images side by side
                                left = Inches(0.5 + (images_added * 4.5))
                                max_width = Inches(4)
                                max_height = Inches(3)
                            else:
                                # Multiple images - smaller grid
                                left = Inches(0.5 + (images_added % 2) * 4.5)
                                if images_added >= 2:
                                    y_pos = Inches(4.0)
                                max_width = Inches(4)
                                max_height = Inches(2.5)
                            
                            # Add image with preserved aspect ratio
                            pic = slide.shapes.add_picture(img_path, left, y_pos)
                            
                            # Scale to fit within max dimensions while preserving aspect ratio
                            if pic.width > max_width:
                                ratio = max_width / pic.width
                                pic.width = max_width
                                pic.height = int(pic.height * ratio)
                            
                            if pic.height > max_height:
                                ratio = max_height / pic.height
                                pic.height = max_height
                                pic.width = int(pic.width * ratio)
                            
                            images_added += 1
                            
                            # Add image caption
                            caption_shape = slide.shapes.add_textbox(left, y_pos + pic.height, pic.width, Inches(0.3))
                            caption_frame = caption_shape.text_frame
                            caption_para = caption_frame.paragraphs[0]
                            caption_para.text = img_file.replace('.png', '').replace('_', ' ').title()
                            caption_para.font.size = Pt(10)
                            caption_para.alignment = PP_ALIGN.CENTER
                            
                        except Exception as e:
                            print(f"⚠️  Could not add image {img_file}: {e}")
    
    # Save presentation
    try:
        prs.save(presentation_path)
        print(f"✅ PowerPoint saved: {presentation_path}")
        return True
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(
        description='Run fMRI QA analysis with PowerPoint report generation',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python run_qa_with_ppt_V2_simple.py --func_dir /path/to/func/
  python run_qa_with_ppt_V2_simple.py --func_dir /path/to/func/ --output_name "MyQA_Report.pptx"
  python run_qa_with_ppt_V2_simple.py --func_dir /path/to/func/ --pattern "*bold*" --no_ppt
        """
    )
    
    parser.add_argument('--func_dir', required=True,
                        help='Path to func directory with BOLD files')
    parser.add_argument('--pattern', default='*fmri*',
                        help='Filename pattern (default: *fmri*)')
    parser.add_argument('--extension', default='.nii.gz',
                        help='File extension (default: .nii.gz)')
    parser.add_argument('--output_name', 
                        help='Custom PowerPoint filename')
    parser.add_argument('--no_ppt', action='store_true',
                        help='Skip PowerPoint generation')
    
    args = parser.parse_args()
    
    # First run the existing QA analysis
    print("🚀 Running QA analysis...")
    
    # Use the existing run_qa_simple.py script
    qa_script = os.path.join(os.path.dirname(__file__), 'run_qa_simple.py')
    cmd = [
        sys.executable, qa_script,
        '--func_dir', args.func_dir,
        '--pattern', args.pattern
    ]
    
    try:
        result = subprocess.run(cmd, check=True, capture_output=True, text=True)
        print("✅ QA analysis completed")
        print(result.stdout)
    except subprocess.CalledProcessError as e:
        print(f"❌ QA analysis failed: {e}")
        print(e.stderr)
        return 1
    
    # Find the generated QA output directories
    parent_dir = os.path.dirname(args.func_dir)
    qa_dirs = [d for d in os.listdir(parent_dir) if d.startswith('qa_output_') and os.path.isdir(os.path.join(parent_dir, d))]
    qa_output_dirs = [os.path.join(parent_dir, d) for d in qa_dirs]
    
    if not qa_output_dirs:
        print("❌ No QA output directories found")
        return 1
    
    print(f"📁 Found {len(qa_output_dirs)} QA output directories")
    
    # Create PowerPoint presentation
    if not args.no_ppt and PPTX_AVAILABLE:
        if args.output_name is None:
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            args.output_name = f"fMRI_QA_Report_{timestamp}.pptx"
        
        ppt_path = os.path.join(parent_dir, args.output_name)
        
        # Extract metrics (simplified)
        summary_data = extract_metrics_from_qa_dirs(qa_output_dirs)
        
        print(f"\n🎨 Creating PowerPoint presentation...")
        ppt_success = create_qa_powerpoint(qa_output_dirs, ppt_path, summary_data)
        
        if ppt_success:
            print(f"📊 PowerPoint report created: {ppt_path}")
        else:
            print(f"❌ Failed to create PowerPoint report")
    elif not PPTX_AVAILABLE:
        print(f"\n⚠️  PowerPoint generation skipped - install python-pptx to enable")
    
    return 0

if __name__ == "__main__":
    # Check for dependencies
    if not PPTX_AVAILABLE:
        print("Installing python-pptx for PowerPoint generation...")
        try:
            import subprocess
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
            print("✅ python-pptx installed successfully!")
            # Restart to import the module
            os.execv(sys.executable, ['python'] + sys.argv)
        except Exception as e:
            print(f"❌ Failed to install python-pptx: {e}")
            print("Please install manually: pip install python-pptx")
            sys.exit(1)
    
    sys.exit(main())