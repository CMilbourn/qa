#!/usr/bin/env python3
# 20260527
# fsleyes_nii_to_pptx.py
# example run line /opt/homebrew/bin/python3 /Users/cmilbourn/Documents/GitHub/qa/fMRI_report_python/fsleyes_nii_to_pptx.py \
# --input-dir /Users/cmilbourn/Documents/Sweet_Data/Development_Data_Phase3/nifti/20260620_phantomrun_5/Run1/ --output-pptx /Users/cmilbourn/Documents/tSNR_check_40dyn/20260619_fsleyes_review.pptx --recursive --keep-png

import argparse
import shutil
import subprocess
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
except ImportError as exc:
    print("python-pptx is required. Install with: pip install python-pptx")
    raise SystemExit(1) from exc


def find_nii_files(input_dir: Path, recursive: bool) -> list[Path]:
    pattern = "**/*.nii.gz" if recursive else "*.nii.gz"
    files = sorted(input_dir.glob(pattern))
    filtered = []
    for f in files:
        if not f.is_file():
            continue
        name_l = f.name.lower()
        if "_brain" in name_l:
            continue
        filtered.append(f)
    return filtered


def run_fsleyes_render(nii_file: Path, png_file: Path, scene: str, width: int, height: int, fsleyes_bin: str) -> None:
    cmd = [
        fsleyes_bin,
        "render",
        "--hideCursor",
        "--scene",
        scene,
        "--size",
        str(width),
        str(height),
        "--outfile",
        str(png_file),
        str(nii_file),
        "-cm",
        "hot",
    ]

    proc = subprocess.run(cmd, capture_output=True, text=True)
    if proc.returncode != 0:
        stderr = proc.stderr.strip()
        stdout = proc.stdout.strip()
        msg = stderr if stderr else stdout
        raise RuntimeError(f"FSLeyes render failed for {nii_file.name}: {msg}")



def add_slide_with_image(prs: Presentation, title_text: str, image_path: Path) -> None:
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.5))
    tf = title_box.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    # Fit image into a 16:9 slide area under title.
    max_left = Inches(0.4)
    max_top = Inches(0.8)
    max_width = Inches(12.5)
    max_height = Inches(6.2)

    from PIL import Image

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    # Assume 96 DPI for layout conversion.
    img_w_in = img_w / 96.0
    img_h_in = img_h / 96.0

    width_ratio = float(max_width) / float(Inches(img_w_in))
    height_ratio = float(max_height) / float(Inches(img_h_in))
    scale_ratio = min(width_ratio, height_ratio, 1.0)

    final_w = Inches(img_w_in) * scale_ratio
    final_h = Inches(img_h_in) * scale_ratio

    left = max_left + (max_width - final_w) / 2
    top = max_top + (max_height - final_h) / 2

    slide.shapes.add_picture(str(image_path), left, top, width=final_w, height=final_h)



def build_pptx(rendered_pngs: list[tuple[Path, Path]], output_pptx: Path, title: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = f"Total images: {len(rendered_pngs)}"

    for nii_file, png_file in rendered_pngs:
        add_slide_with_image(prs, nii_file.name, png_file)

    prs.save(str(output_pptx))



def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Render .nii.gz files with FSLeyes and create a PowerPoint with one image per slide."
    )
    parser.add_argument("--input-dir", required=True, help="Folder containing .nii.gz files")
    parser.add_argument("--output-pptx", required=True, help="Output PowerPoint file path")
    parser.add_argument("--render-dir", default=None, help="Optional folder for rendered PNG files")
    parser.add_argument("--title", default="FSLeyes NIfTI Review", help="Title for the first slide")
    parser.add_argument("--scene", default="ortho", choices=["ortho", "lightbox", "3d"], help="FSLeyes scene type")
    parser.add_argument("--width", type=int, default=1600, help="Rendered PNG width in pixels")
    parser.add_argument("--height", type=int, default=1200, help="Rendered PNG height in pixels")
    parser.add_argument("--recursive", action="store_true", help="Search input directory recursively")
    parser.add_argument("--keep-png", action="store_true", help="Keep rendered PNGs after PPTX is created")
    parser.add_argument("--fsleyes-bin", default="fsleyes", help="Path to fsleyes executable")
    return parser.parse_args()



def main() -> int:
    args = parse_args()

    input_dir = Path(args.input_dir).expanduser().resolve()
    output_pptx = Path(args.output_pptx).expanduser().resolve()

    if not input_dir.is_dir():
        print(f"Input directory does not exist: {input_dir}")
        return 1

    if shutil.which(args.fsleyes_bin) is None:
        print(f"Could not find fsleyes executable: {args.fsleyes_bin}")
        print("Install FSLeyes or provide --fsleyes-bin with a full path.")
        return 1

    nii_files = find_nii_files(input_dir, recursive=args.recursive)
    if not nii_files:
        print(f"No .nii.gz files found in: {input_dir}")
        return 1

    if args.render_dir:
        render_dir = Path(args.render_dir).expanduser().resolve()
    else:
        render_dir = output_pptx.parent / (output_pptx.stem + "_renders")

    render_dir.mkdir(parents=True, exist_ok=True)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    rendered = []
    failures = []

    print(f"Found {len(nii_files)} NIfTI files")
    for idx, nii_file in enumerate(nii_files, start=1):
        png_name = nii_file.name.replace(".nii.gz", ".png")
        png_file = render_dir / png_name

        print(f"[{idx}/{len(nii_files)}] Rendering {nii_file.name}")
        try:
            run_fsleyes_render(
                nii_file=nii_file,
                png_file=png_file,
                scene=args.scene,
                width=args.width,
                height=args.height,
                fsleyes_bin=args.fsleyes_bin,
            )
            rendered.append((nii_file, png_file))
        except Exception as exc:
            failures.append((nii_file, str(exc)))
            print(f"  FAILED: {exc}")

    if not rendered:
        print("No images were rendered successfully; skipping PowerPoint creation.")
        return 1

    build_pptx(rendered, output_pptx, args.title)
    print(f"PowerPoint created: {output_pptx}")

    if failures:
        print("\nRender failures:")
        for nii_file, err in failures:
            print(f"  - {nii_file.name}: {err}")

    if not args.keep_png and not args.render_dir:
        try:
            shutil.rmtree(render_dir)
            print(f"Removed temporary render directory: {render_dir}")
        except Exception as exc:
            print(f"Could not remove temporary render directory {render_dir}: {exc}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
