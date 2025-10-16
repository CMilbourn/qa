#!/usr/bin/env python3
"""
Simple PowerPoint Creator for sub003 QA Results
Version: 1.0.0

Creates a comprehensive PowerPoint presentation from all QA output directories.
Adapted for sub003 data structure which uses original QA script outputs.

Version History:
- v1.0.0 (2025-10-16): Baseline version adapted from sub001 for sub003 data structure
"""

import os
import sys
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, content=""):
    """Add a content slide with title."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if content:
        slide.placeholders[1].text = content
    
    return slide

def add_image_slide(prs, title, image_path, description=""):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add image (preserve aspect ratio)
    if os.path.exists(image_path):
        try:
            # Calculate position to center the image
            img_left = Inches(1)
            img_top = Inches(1.2)
            img_width = Inches(8)
            
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    
    # Add description if provided
    if description:
        desc_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        desc_frame = desc_shape.text_frame
        desc_frame.text = description
        desc_paragraph = desc_frame.paragraphs[0]
        desc_paragraph.font.size = Pt(12)
        desc_paragraph.alignment = PP_ALIGN.CENTER

def extract_tr_from_json(qa_parent_dir, qa_dir_name):
    """Extract TR from JSON metadata file for sub003."""
    # Remove qa_output_ prefix to get NIFTI filename
    nifti_name = qa_dir_name.replace('qa_output_', '')
    json_file = os.path.join(qa_parent_dir, f"{nifti_name}.json")
    
    if os.path.exists(json_file):
        try:
            with open(json_file, 'r') as f:
                metadata = json.load(f)
                return metadata.get('RepetitionTime', 'Unknown')
        except Exception as e:
            print(f"Error reading {json_file}: {e}")
    return 'Unknown'

def get_scan_info(qa_name):
    """Extract scan information from sub003 QA directory name."""
    if "2mm_longerTR" in qa_name:
        return "2mm Resolution (Long TR)", "2mm × 2mm × 2mm (Long TR)"
    elif "3mm" in qa_name:
        return "3mm Resolution", "3mm × 3mm × 3mm"
    elif "2mm_pre" in qa_name and "de-9" in qa_name:
        return "2mm Resolution (Pre-scan 1)", "2mm × 2mm × 2mm (Pre-scan 1)"
    elif "2mm_pre" in qa_name and "de-2" in qa_name:
        return "2mm Resolution (Pre-scan 2)", "2mm × 2mm × 2mm (Pre-scan 2)"
    else:
        return "fMRI Acquisition", "Unknown Resolution"

def count_qa_images(qa_dir):
    """Count QA images in directory."""
    try:
        png_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
        tsnr_files = [f for f in png_files if 'tsnr' in f.lower() or 'tSNR' in f]
        return len(png_files), len(tsnr_files)
    except Exception as e:
        print(f"Error counting images in {qa_dir}: {e}")
        return 0, 0

def create_comprehensive_ppt(qa_parent_dir, output_file):
    """Create comprehensive PowerPoint from all sub003 QA directories."""
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, "fMRI Quality Assessment Report", "sub003 Complete Analysis")
    
    # Get all QA directories
    qa_dirs = []
    if os.path.exists(qa_parent_dir):
        for item in os.listdir(qa_parent_dir):
            item_path = os.path.join(qa_parent_dir, item)
            if os.path.isdir(item_path) and item.startswith('qa_output_'):
                qa_dirs.append((item, item_path))
    
    qa_dirs.sort()  # Sort alphabetically
    
    if not qa_dirs:
        print(f"❌ No QA directories found in {qa_parent_dir}")
        return False
    
    print(f"📊 Found {len(qa_dirs)} QA result directories")
    
    # Summary slide
    summary_slide = add_content_slide(prs, "Processing Summary", "")
    
    summary_text = f"Total scans processed: {len(qa_dirs)}\\n\\n"
    
    for qa_name, qa_dir in qa_dirs:
        scan_title, scan_desc = get_scan_info(qa_name)
        tr = extract_tr_from_json(qa_parent_dir, qa_name)
        total_images, tsnr_images = count_qa_images(qa_dir)
        
        summary_text += f"• {scan_title}\\n"
        summary_text += f"  TR: {tr}s, QA Images: {total_images}, tSNR Images: {tsnr_images}\\n\\n"
    
    summary_slide.placeholders[1].text = summary_text
    
    # Add slides for each QA directory
    # Use sub003-specific image file names
    image_types = [
        ('Mean_image.png', 'Mean Signal Image'),
        ('mean_montage.png', 'Mean Signal Montage'),
        ('tSNR_montage.png', 'Temporal SNR Montage (Orange)'),
        ('tSNR_w_ROI_images.png', 'tSNR with ROI Analysis'),
        ('tSNR_per_unit_time.png', 'tSNR per Unit Time'),
        ('isnr_montage.png', 'Instantaneous SNR Montage'),
        ('iSNR_cor.png', 'iSNR Coronal View'),
        ('iSNR_sag.png', 'iSNR Sagittal View'),
        ('masked_noise.png', 'Masked Noise Analysis'),
        ('masked_noise_volume_montage.png', 'Noise Volume Montage'),
        ('TS_images.png', 'Time Series Analysis')
    ]
    
    for qa_name, qa_dir in qa_dirs:
        print(f"Processing: {qa_name}")
        
        # Extract scan info
        scan_title, scan_desc = get_scan_info(qa_name)
        tr = extract_tr_from_json(qa_parent_dir, qa_name)
        total_images, tsnr_images = count_qa_images(qa_dir)
        
        # Section title slide
        subtitle = f"TR: {tr}s | {scan_desc} | QA Images: {total_images} | tSNR Images: {tsnr_images}"
        add_title_slide(prs, scan_title, subtitle)
        
        # Add image slides
        for image_file, image_title in image_types:
            image_path = os.path.join(qa_dir, image_file)
            
            if os.path.exists(image_path):
                full_title = f"{scan_title}: {image_title}"
                
                # Add special descriptions for key images
                description = ""
                if 'tSNR_montage' in image_file:
                    description = "Orange/Red colors indicate higher temporal SNR values. Brighter regions show better signal stability over time."
                elif 'tSNR_w_ROI' in image_file:
                    description = "Temporal SNR analysis with Region of Interest overlay showing quantitative measurements."
                elif 'tSNR_per_unit_time' in image_file:
                    description = f"Temporal SNR normalized per unit time (TR={tr}s). Accounts for different repetition times."
                elif 'isnr_montage' in image_file:
                    description = "Instantaneous SNR montage showing spatial distribution of signal-to-noise ratio."
                
                add_image_slide(prs, full_title, image_path, description)
                print(f"  Added: {image_title}")
            else:
                print(f"  ⚠️  Missing: {image_file}")
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✅ PowerPoint saved: {output_file}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        print(f"📊 File size: {file_size:.1f} MB")
        print(f"📄 Total slides: {len(prs.slides)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    """Main function."""
    if len(sys.argv) < 3:
        print("Usage: python create_sub003_ppt.py <qa_parent_dir> <output_file>")
        print("Example: python create_sub003_ppt.py /path/to/sub003/func/ sub003_QA_report.pptx")
        sys.exit(1)
    
    qa_parent_dir = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(qa_parent_dir):
        print(f"❌ QA directory not found: {qa_parent_dir}")
        sys.exit(1)
    
    print(f"🚀 Creating sub003 PowerPoint report from: {qa_parent_dir}")
    
    success = create_comprehensive_ppt(qa_parent_dir, output_file)
    
    if success:
        print(f"🎉 sub003 Report generation complete!")
    else:
        print(f"❌ Report generation failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()