#!/usr/bin/env python
"""
Create PowerPoint from existing QA output directories
Version: 1.0.0
Usage: python create_ppt_from_qa.py --qa_parent_dir /path/to/parent/ --output_name "report.pptx"

Version History:
- v1.0.0 (2025-10-16): Baseline version with comprehensive QA reporting
"""

import argparse
import os
import json
from datetime import datetime

# PowerPoint libraries
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed.")

def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception as e:
        print(f"Warning: Could not read TR from JSON: {e}")
    return None

def extract_metrics_from_qa_dirs(qa_output_dirs, func_dir=None):
    """Extract metrics from existing QA output directories and original data"""
    metrics_list = []
    
    for qa_dir in qa_output_dirs:
        if not os.path.exists(qa_dir):
            continue
            
        # Extract filename from directory name
        dirname = os.path.basename(qa_dir)
        if dirname.startswith('qa_output_'):
            filename = dirname[10:]  # Remove 'qa_output_' prefix
        else:
            filename = dirname
            
        # Try to get TR from func directory
        tr_value = None
        ernst_factor = None
        data_shape = None
        
        if func_dir:
            for ext in ['.nii.gz', '.nii']:
                nifti_file = os.path.join(func_dir, filename + ext)
                if os.path.exists(nifti_file):
                    tr_value = get_tr_from_json(nifti_file)
                    if tr_value:
                        # Calculate Ernst scaling factor
                        if tr_value <= 0.7:
                            ernst_factor = 0.5745  # sin(35°)
                        elif tr_value <= 1.0:
                            ernst_factor = 0.7071  # sin(45°)
                        elif tr_value <= 1.5:
                            ernst_factor = 0.8155  # sin(54.7°)
                        else:
                            ernst_factor = 1.0     # sin(90°)
                    
                    # Try to get data shape from NIFTI file
                    try:
                        import nibabel as nib
                        img = nib.load(nifti_file)
                        data_shape = img.shape
                    except Exception as e:
                        print(f"Could not load NIFTI file {nifti_file}: {e}")
                    break
        
        # Try to extract metrics from NIFTI files in QA directory
        isnr_value = None
        tsnr_value = None
        mean_tsnr_unit_time = None
        
        # Look for iSNR and tSNR NIFTI files
        isnr_file = os.path.join(qa_dir, 'isnr.nii.gz')
        tsnr_file = os.path.join(qa_dir, 'tsnr.nii.gz')
        
        try:
            import nibabel as nib
            import numpy as np
            
            if os.path.exists(isnr_file):
                isnr_img = nib.load(isnr_file)
                isnr_data = isnr_img.get_fdata()
                isnr_value = float(np.mean(isnr_data[isnr_data > 0]))  # Mean of non-zero voxels
            
            if os.path.exists(tsnr_file):
                tsnr_img = nib.load(tsnr_file)
                tsnr_data = tsnr_img.get_fdata()
                tsnr_value = float(np.mean(tsnr_data[tsnr_data > 0]))  # Mean of non-zero voxels
                
                # Calculate tSNR per unit time if TR is available
                if tr_value and tsnr_value:
                    mean_tsnr_unit_time = tsnr_value / np.sqrt(tr_value)
        
        except Exception as e:
            print(f"Could not extract metrics from NIFTI files in {qa_dir}: {e}")
        
        metrics = {
            'filename': filename,
            'tr': tr_value if tr_value else 'Unknown',
            'ernst_factor': ernst_factor,
            'data_shape': data_shape,
            'isnr': isnr_value,
            'tsnr': tsnr_value,
            'tsnr_per_unit_time': mean_tsnr_unit_time,
            'qa_dir': qa_dir
        }
        
        metrics_list.append(metrics)
    
    return metrics_list

def create_qa_powerpoint(qa_output_dirs, presentation_path, summary_data):
    """Create a PowerPoint presentation with QA results"""
    if not PPTX_AVAILABLE:
        print("❌ Cannot create PowerPoint - python-pptx not installed")
        return False
    
    print("🎨 Creating PowerPoint presentation...")
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assessment Report"
    subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n{len(qa_output_dirs)} datasets analyzed"
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "QA Summary"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = "Dataset Summary:"
    p.font.bold = True
    p.font.size = Pt(16)
    
    for i, data in enumerate(summary_data):
        p = tf.add_paragraph()
        p.text = f"\n{i+1}. {data['filename'][:50]}{'...' if len(data['filename']) > 50 else ''}"
        p.font.size = Pt(12)
        p.font.bold = True
        
        # Build summary metrics line
        metrics_line = f"   • TR: {data['tr']}s"
        if data['isnr']:
            metrics_line += f" | iSNR: {data['isnr']:.2f}"
        if data['tsnr']:
            metrics_line += f" | tSNR: {data['tsnr']:.2f}"
        if data['ernst_factor']:
            metrics_line += f" | Ernst: {data['ernst_factor']:.3f}"
        
        p = tf.add_paragraph()
        p.text = metrics_line
        p.font.size = Pt(10)
        
        qa_dir = data['qa_dir']
        if os.path.exists(qa_dir):
            image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
            shape_text = ""
            if data['data_shape'] and len(data['data_shape']) >= 3:
                shape_text = f" | Shape: {data['data_shape'][0]}×{data['data_shape'][1]}×{data['data_shape'][2]}"
            
            p = tf.add_paragraph()
            p.text = f"   • {len(image_files)} QA images{shape_text}"
            p.font.size = Pt(10)
    
    # Create slides for each dataset
    for i, (qa_dir, data) in enumerate(zip(qa_output_dirs, summary_data)):
        
        if not os.path.exists(qa_dir):
            continue
            
        # Dataset overview slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Dataset {i+1}: {data['filename']}"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = f"File: {data['filename']}"
        p.font.size = Pt(14)
        p.font.bold = True
        
        # Create comprehensive metrics text
        image_files = [f for f in os.listdir(qa_dir) if f.endswith('.png')]
        
        info_text = f"\nQA Metrics Summary:\n"
        
        # Add TR and Ernst scaling info
        if data['tr'] != 'Unknown':
            info_text += f"• TR (Repetition Time): {data['tr']}s\n"
            if data['ernst_factor']:
                info_text += f"• Ernst scaling factor: {data['ernst_factor']:.4f}\n"
        
        # Add data shape
        if data['data_shape']:
            if len(data['data_shape']) == 4:
                info_text += f"• Image data shape: {data['data_shape'][0]} × {data['data_shape'][1]} × {data['data_shape'][2]} × {data['data_shape'][3]} voxels\n"
            else:
                info_text += f"• Image data shape: {data['data_shape']}\n"
        
        # Add SNR metrics
        if data['isnr']:
            info_text += f"• Image SNR (iSNR): {data['isnr']:.2f}\n"
        
        if data['tsnr']:
            info_text += f"• Temporal SNR (tSNR): {data['tsnr']:.2f}\n"
        
        if data['tsnr_per_unit_time']:
            info_text += f"• tSNR per unit time: {data['tsnr_per_unit_time']:.2f}\n"
        
        info_text += f"\nTechnical Details:\n"
        info_text += f"• QA output directory: {os.path.basename(qa_dir)}\n"
        info_text += f"• Generated QA images: {len(image_files)} files\n"
        
        # List key image types
        key_images = []
        for img_file in sorted(image_files):
            if any(key in img_file.lower() for key in ['isnr', 'tsnr', 'mean', 'montage']):
                key_images.append(img_file)
        
        if key_images:
            info_text += f"\nKey QA Images:\n"
            for img_file in key_images[:6]:  # Limit to avoid text overflow
                info_text += f"• {img_file}\n"
            if len(key_images) > 6:
                info_text += f"• ... and {len(key_images) - 6} more images\n"
        
        p = tf.add_paragraph()
        p.text = info_text
        p.font.size = Pt(12)
        
        # Group and display images
        image_groups = {
            'Mean Images': [f for f in image_files if 'mean' in f.lower()],
            'SNR Analysis': [f for f in image_files if any(x in f.lower() for x in ['isnr', 'tsnr', 'snr'])],
            'Montage Views': [f for f in image_files if 'montage' in f.lower()],
            'Noise Analysis': [f for f in image_files if 'noise' in f.lower()],
            'Time Series': [f for f in image_files if any(x in f.lower() for x in ['ts_', 'timeseries'])],
            'Other': [f for f in image_files if not any(keyword in f.lower() for keyword in ['mean', 'snr', 'montage', 'noise', 'ts_', 'timeseries'])]
        }
        
        for category, images in image_groups.items():
            if not images:
                continue
                
            # Create slide for this category
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            title_frame = title_shape.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = f"Dataset {i+1}: {category}"
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.alignment = PP_ALIGN.CENTER
            
            # Add images (up to 4 per slide)
            y_pos = Inches(1.0)
            for j, img_file in enumerate(images[:4]):
                img_path = os.path.join(qa_dir, img_file)
                if os.path.exists(img_path):
                    try:
                        # Position images in a 2x2 grid with preserved aspect ratio
                        left = Inches(0.5 + (j % 2) * 4.75)
                        if j >= 2:
                            y_pos = Inches(4.0)
                        
                        # Add image with preserved aspect ratio (no width/height specified)
                        pic = slide.shapes.add_picture(img_path, left, y_pos)
                        
                        # Scale down if too large, maintaining aspect ratio
                        max_width = Inches(4.5)
                        max_height = Inches(3)
                        
                        if pic.width > max_width:
                            ratio = max_width / pic.width
                            pic.width = max_width
                            pic.height = int(pic.height * ratio)
                        
                        if pic.height > max_height:
                            ratio = max_height / pic.height
                            pic.height = max_height
                            pic.width = int(pic.width * ratio)
                        
                        # Add caption below the actual image
                        caption_shape = slide.shapes.add_textbox(left, y_pos + pic.height, pic.width, Inches(0.5))
                        caption_frame = caption_shape.text_frame
                        caption_para = caption_frame.paragraphs[0]
                        caption_para.text = img_file.replace('.png', '').replace('_', ' ').title()
                        caption_para.font.size = Pt(9)
                        caption_para.alignment = PP_ALIGN.CENTER
                        
                    except Exception as e:
                        print(f"⚠️  Could not add image {img_file}: {e}")
    
    # Save presentation
    try:
        prs.save(presentation_path)
        print(f"✅ PowerPoint saved: {presentation_path}")
        return True
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(description='Create PowerPoint from existing QA output directories')
    parser.add_argument('--qa_parent_dir', required=True, help='Parent directory containing qa_output_* folders')
    parser.add_argument('--func_dir', help='Optional func directory to extract TR values')
    parser.add_argument('--output_name', default='QA_Report.pptx', help='PowerPoint filename')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.qa_parent_dir):
        print(f"❌ Directory does not exist: {args.qa_parent_dir}")
        return 1
    
    # Find QA output directories
    qa_dirs = []
    for item in os.listdir(args.qa_parent_dir):
        if item.startswith('qa_output_') and os.path.isdir(os.path.join(args.qa_parent_dir, item)):
            qa_dirs.append(os.path.join(args.qa_parent_dir, item))
    
    if not qa_dirs:
        print(f"❌ No qa_output_* directories found in {args.qa_parent_dir}")
        return 1
    
    print(f"📁 Found {len(qa_dirs)} QA output directories:")
    for qa_dir in qa_dirs:
        print(f"   • {os.path.basename(qa_dir)}")
    
    # Extract metrics
    summary_data = extract_metrics_from_qa_dirs(qa_dirs, args.func_dir)
    
    # Create PowerPoint
    ppt_path = os.path.join(args.qa_parent_dir, args.output_name)
    success = create_qa_powerpoint(qa_dirs, ppt_path, summary_data)
    
    if success:
        print(f"🎉 PowerPoint report created successfully!")
        print(f"📄 File: {ppt_path}")
    else:
        print("❌ Failed to create PowerPoint report")
    
    return 0 if success else 1

if __name__ == "__main__":
    exit(main())