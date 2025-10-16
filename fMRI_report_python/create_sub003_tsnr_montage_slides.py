#!/usr/bin/env python3
"""
tSNR Montage Slide Creator for sub003 QA Results
Version: 1.3.0

Creates PowerPoint slides focused specifically on tSNR montage visualizations
with the orange colormap format showing all 55 slices as in the reference slide.

Version History:
- v1.0.0 (2025-10-16): Initial version creating tSNR montage slides for sub003
- v1.1.0 (2025-10-16): Fixed to use correct tSNR images with all 55 slices and proper orange colormap
- v1.2.0 (2025-10-16): Updated to use tSNR_montage.png (2MB file) with proper orange colormap
- v1.3.0 (2025-10-16): Added tSNR metric extraction and display on slides

Usage: python create_sub003_tsnr_montage_slides.py <qa_parent_dir> <output_file>
"""

import os
import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_tsnr_montage_slide(prs, title, image_path, tr_value="Unknown", ernst_scaling=1.0, tsnr_metrics=None):
    """Add a slide with tSNR montage image matching the reference format."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title at top (matching reference format)
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(28)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add subtitle with Ernst scaling info (matching reference)
    subtitle_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
    subtitle_frame = subtitle_shape.text_frame
    subtitle_frame.text = f"Temporal SNR (Ernst scaled by {ernst_scaling:.3f})"
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(16)
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add tSNR montage image (centered, preserve aspect ratio)
    if os.path.exists(image_path):
        try:
            img_left = Inches(1)
            img_top = Inches(1.5)
            img_width = Inches(8)
            
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            
        except Exception as e:
            print(f"⚠️  Error adding tSNR montage {image_path}: {e}")
            # Add error text if image fails
            error_shape = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
            error_frame = error_shape.text_frame
            error_frame.text = f"⚠️ tSNR montage image not found:\n{os.path.basename(image_path)}"
            error_paragraph = error_frame.paragraphs[0]
            error_paragraph.font.size = Pt(14)
            error_paragraph.alignment = PP_ALIGN.CENTER
    else:
        # Add missing image text
        missing_shape = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1))
        missing_frame = missing_shape.text_frame
        missing_frame.text = f"❌ tSNR montage image missing:\n{os.path.basename(image_path)}"
        missing_paragraph = missing_frame.paragraphs[0]
        missing_paragraph.font.size = Pt(14)
        missing_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add tSNR metrics at bottom if available
    if tsnr_metrics:
        metrics_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(1))
        metrics_frame = metrics_shape.text_frame
        
        metrics_text = "tSNR Metrics: "
        if 'mean_tsnr' in tsnr_metrics:
            metrics_text += f"Mean: {tsnr_metrics['mean_tsnr']:.2f}  "
        if 'median_tsnr' in tsnr_metrics:
            metrics_text += f"Median: {tsnr_metrics['median_tsnr']:.2f}  "
        if 'std_tsnr' in tsnr_metrics:
            metrics_text += f"Std: {tsnr_metrics['std_tsnr']:.2f}  "
        if 'max_tsnr' in tsnr_metrics:
            metrics_text += f"Max: {tsnr_metrics['max_tsnr']:.2f}"
        
        metrics_frame.text = metrics_text
        metrics_paragraph = metrics_frame.paragraphs[0]
        metrics_paragraph.font.size = Pt(14)
        metrics_paragraph.font.bold = True
        metrics_paragraph.alignment = PP_ALIGN.CENTER

def extract_tr_from_json(qa_dir):
    """Extract TR from JSON metadata file for sub003."""
    # Look for JSON files with TR information
    for file in os.listdir(qa_dir):
        if file.endswith('.json') and 'bold' in file.lower():
            json_path = os.path.join(qa_dir, file)
            try:
                with open(json_path, 'r') as f:
                    metadata = json.load(f)
                    if 'RepetitionTime' in metadata:
                        return metadata['RepetitionTime']
            except:
                continue
    return "Unknown"

def get_scan_info(qa_name):
    """Extract scan information from sub003 QA directory name."""
    # Parse the qa_output directory name to extract scan info
    if "2mm_pre" in qa_name:
        return "2mm Resolution Scan"
    elif "2mm_longerTR" in qa_name:
        return "2mm Resolution Scan (Longer TR)"
    elif "3mm" in qa_name:
        return "3mm Resolution Scan"
    elif "2mm_1.5" in qa_name:
        return "2mm x 1.5mm Resolution Scan"
    elif "1.5_mm_iso" in qa_name:
        return "1.5mm Isotropic Resolution Scan"
    else:
        return "fMRI Scan"

def calculate_ernst_scaling(tr_value):
    """Calculate Ernst angle scaling factor."""
    try:
        tr_float = float(tr_value)
        # Ernst angle calculation (simplified)
        if tr_float > 0:
            return 1.0  # Placeholder - you can implement actual Ernst scaling if needed
    except:
        pass
    return 1.0

def extract_tsnr_metrics(qa_dir):
    """Extract tSNR metrics from QA directory output files."""
    tsnr_metrics = {}
    
    # Look for text files with tSNR metrics
    potential_files = [
        'qa_metrics.txt',
        'tsnr_metrics.txt',
        'summary.txt',
        'results.txt'
    ]
    
    for filename in potential_files:
        file_path = os.path.join(qa_dir, filename)
        if os.path.exists(file_path):
            try:
                with open(file_path, 'r') as f:
                    content = f.read().lower()
                    
                    # Extract tSNR values using simple parsing
                    import re
                    
                    # Look for patterns like "mean tsnr: 45.67" or "tsnr mean: 45.67"
                    mean_match = re.search(r'(?:mean\s+tsnr|tsnr\s+mean):\s*([0-9]+\.?[0-9]*)', content)
                    if mean_match:
                        tsnr_metrics['mean_tsnr'] = float(mean_match.group(1))
                    
                    # Look for median tSNR
                    median_match = re.search(r'(?:median\s+tsnr|tsnr\s+median):\s*([0-9]+\.?[0-9]*)', content)
                    if median_match:
                        tsnr_metrics['median_tsnr'] = float(median_match.group(1))
                    
                    # Look for std tSNR
                    std_match = re.search(r'(?:std\s+tsnr|tsnr\s+std|standard\s+deviation\s+tsnr):\s*([0-9]+\.?[0-9]*)', content)
                    if std_match:
                        tsnr_metrics['std_tsnr'] = float(std_match.group(1))
                    
                    # Look for max tSNR
                    max_match = re.search(r'(?:max\s+tsnr|tsnr\s+max|maximum\s+tsnr):\s*([0-9]+\.?[0-9]*)', content)
                    if max_match:
                        tsnr_metrics['max_tsnr'] = float(max_match.group(1))
                        
            except Exception as e:
                print(f"  ⚠️ Could not read {filename}: {e}")
                continue
    
    # If no text files found, try to extract from NIfTI files using nibabel
    if not tsnr_metrics:
        nifti_file = os.path.join(qa_dir, 'tsnr_tsnr_map.nii.gz')
        if os.path.exists(nifti_file):
            try:
                import nibabel as nib
                import numpy as np
                
                print(f"  📊 Computing tSNR metrics from NIfTI file...")
                tsnr_img = nib.load(nifti_file)
                tsnr_data = tsnr_img.get_fdata()
                
                # Remove zeros and invalid values
                valid_tsnr = tsnr_data[tsnr_data > 0]
                if len(valid_tsnr) > 0:
                    tsnr_metrics['mean_tsnr'] = float(np.mean(valid_tsnr))
                    tsnr_metrics['median_tsnr'] = float(np.median(valid_tsnr))
                    tsnr_metrics['std_tsnr'] = float(np.std(valid_tsnr))
                    tsnr_metrics['max_tsnr'] = float(np.max(valid_tsnr))
                    
            except Exception as e:
                print(f"  ⚠️ Could not compute tSNR metrics from NIfTI: {e}")
    
    return tsnr_metrics if tsnr_metrics else None

def create_sub003_tsnr_slides(qa_parent_dir, output_file):
    """Create PowerPoint with tSNR montage slides for sub003."""
    
    # Create presentation
    prs = Presentation()
    
    # Get all QA directories
    qa_dirs = []
    if os.path.exists(qa_parent_dir):
        for item in os.listdir(qa_parent_dir):
            item_path = os.path.join(qa_parent_dir, item)
            if os.path.isdir(item_path) and item.startswith('qa_output'):
                qa_dirs.append((item, item_path))
    
    qa_dirs.sort()  # Sort alphabetically
    
    if not qa_dirs:
        print(f"❌ No QA directories found in {qa_parent_dir}")
        return False
    
    print(f"🎨 Creating tSNR montage PowerPoint with {len(qa_dirs)} datasets")
    
    # Title slide
    add_title_slide(prs, "Temporal SNR (tSNR) Montage Analysis", f"sub003 Orange Colormap Visualization - {len(qa_dirs)} Acquisitions")
    
    # Summary slide
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "tSNR Montage Overview"
    
    summary_text = "Temporal SNR Analysis with Orange Colormap:\\n\\n"
    
    for i, (qa_name, qa_dir) in enumerate(qa_dirs, 1):
        scan_type = get_scan_info(qa_name)
        tr = extract_tr_from_json(qa_dir)
        
        summary_text += f"{i}. {scan_type}\\n"
        summary_text += f"   • TR: {tr}s\\n"
        summary_text += f"   • Orange colormap tSNR montage\\n\\n"
    
    slide.placeholders[1].text = summary_text
    
    # Process each QA directory - focus only on tSNR montage
    for qa_name, qa_dir in qa_dirs:
        print(f"Processing: {qa_name}")
        
        scan_type = get_scan_info(qa_name)
        tr = extract_tr_from_json(qa_dir)
        ernst_scaling = calculate_ernst_scaling(tr)
        
        # Extract tSNR metrics
        print(f"  📊 Extracting tSNR metrics...")
        tsnr_metrics = extract_tsnr_metrics(qa_dir)
        if tsnr_metrics:
            print(f"     Mean tSNR: {tsnr_metrics.get('mean_tsnr', 'N/A'):.2f}" if 'mean_tsnr' in tsnr_metrics else "     Mean tSNR: N/A")
        
        # Look for tSNR montage image (orange colormap with all 55 slices)
        tsnr_montage_files = [
            'tSNR_montage.png',       # Large 2MB file with comprehensive montage and orange colormap
            'tsnr_montage.png',       # Alternative case
            'tSNR_w_ROI_images.png',  # Fallback option
            'tSNR_raw.png'            # Last resort
        ]
        
        tsnr_found = False
        for tsnr_file in tsnr_montage_files:
            tsnr_path = os.path.join(qa_dir, tsnr_file)
            if os.path.exists(tsnr_path):
                title = f"{scan_type}: Temporal SNR Montage"
                add_tsnr_montage_slide(prs, title, tsnr_path, tr, ernst_scaling, tsnr_metrics)
                print(f"  ✅ Added: {tsnr_file} with metrics")
                tsnr_found = True
                break
        
        if not tsnr_found:
            print(f"  ⚠️  No tSNR montage image found in {qa_name}")
            # Still create a slide but with missing image placeholder
            title = f"{scan_type}: Temporal SNR Montage"
            add_tsnr_montage_slide(prs, title, "missing_tsnr_montage.png", tr, ernst_scaling, tsnr_metrics)
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✅ tSNR Montage PowerPoint saved: {output_file}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        print(f"📊 File size: {file_size:.1f} MB")
        print(f"📄 Total slides: {len(prs.slides)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    """Main function."""
    if len(sys.argv) < 3:
        print("Usage: python create_sub003_tsnr_montage_slides.py <qa_parent_dir> <output_file>")
        print("Example: python create_sub003_tsnr_montage_slides.py /path/to/sub003/func/ sub003_tSNR_montage.pptx")
        print()
        print("This script creates PowerPoint slides focused on tSNR montage visualizations")
        print("with the orange colormap format, specifically for sub003 data.")
        sys.exit(1)
    
    qa_parent_dir = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(qa_parent_dir):
        print(f"❌ QA directory not found: {qa_parent_dir}")
        sys.exit(1)
    
    print(f"🚀 Creating sub003 tSNR montage slides from: {qa_parent_dir}")
    
    success = create_sub003_tsnr_slides(qa_parent_dir, output_file)
    
    if success:
        print(f"🎉 tSNR Montage report generation complete!")
    else:
        print(f"❌ tSNR Montage report generation failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()