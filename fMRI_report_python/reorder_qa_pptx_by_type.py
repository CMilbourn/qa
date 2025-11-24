#!/usr/bin/env python3
import os
import sys
import re
from glob import glob
from datetime import datetime
from pptx import Presentation

# Robust helpers to inspect slide text

def get_all_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            try:
                txt = shape.text_frame.text or ""
            except Exception:
                txt = ""
            if txt:
                texts.append(txt.strip())
    return "\n".join(texts)


def detect_slide_kind(text: str):
    t = text.strip().lower()
    # Title slide detection
    if "quality assurance report" in t and "combined" in t:
        return "title"
    # Dataset separator slide
    if t.startswith("dataset "):
        return "dataset-sep"

    # Image titles (substring match, lowercase)
    mapping = {
        "Mean Image": ["mean image"],
        "Mean Image Montage - All Slices": ["mean image montage"],
        "Noise Volume Analysis": ["noise volume analysis"],
        "Noise Volume Montage": ["noise volume montage"],
        "Masked Noise Volume Montage": ["masked noise volume montage"],
        "tSNR Map - Sagittal View": ["tsnr map - sagittal"],
        "tSNR Map - Coronal View": ["tsnr map - coronal"],
        "tSNR per Unit Time": ["tsnr per unit time"],
        "Raw tSNR Map": ["raw tsnr map"],
        "tSNR Montage - All Slices": ["tsnr montage"],
        "tSNR with ROI": ["tsnr with roi"],
        "iSNR Map - Sagittal View": ["isnr map - sagittal"],
        "iSNR Map - Coronal View": ["isnr map - coronal"],
        "iSNR Montage - All Slices": ["isnr montage"],
        "Time Series Analysis": ["time series analysis"],
        "Static Spatial Noise": ["static spatial noise"],
    }
    for key, patterns in mapping.items():
        for p in patterns:
            if p in t:
                return key
    return "unknown"


def extract_dataset_index(text: str):
    # Text looks like: "Dataset 1/16\n<name>"
    m = re.search(r"dataset\s+(\d+)\s*/\s*\d+", text.lower())
    if m:
        return int(m.group(1))
    return None


def reorder_by_type(pptx_path: str, output_path: str = None):
    prs = Presentation(pptx_path)

    # Discover slide kinds and dataset indices
    slides_info = []  # (idx, kind, dataset_idx)
    current_dataset = None
    for i, slide in enumerate(prs.slides):
        txt = get_all_text(slide)
        kind = detect_slide_kind(txt)
        if kind == "dataset-sep":
            current_dataset = extract_dataset_index(txt)
        slides_info.append((i, kind, current_dataset))

    # Determine subject order from dataset-sep slides
    dataset_indices = [ds for (_i, k, ds) in slides_info if k == "dataset-sep" and ds is not None]
    dataset_indices = sorted(set(dataset_indices))

    # Desired image ordering (as in the combined report generation)
    type_order = [
        "Mean Image",
        "Mean Image Montage - All Slices",
        "Noise Volume Analysis",
        "Noise Volume Montage",
        "Masked Noise Volume Montage",
        "tSNR Map - Sagittal View",
        "tSNR Map - Coronal View",
        "tSNR per Unit Time",
        "Raw tSNR Map",
        "tSNR Montage - All Slices",
        "tSNR with ROI",
        "iSNR Map - Sagittal View",
        "iSNR Map - Coronal View",
        "iSNR Montage - All Slices",
        "Time Series Analysis",
        "Static Spatial Noise",
    ]

    # Build buckets by kind -> list of slide indices per dataset index (preserving dataset order)
    by_kind = {k: [] for k in type_order}
    title_slide_indices = [i for (i, k, _ds) in slides_info if k == "title"]
    dataset_sep_indices = [i for (i, k, _ds) in slides_info if k == "dataset-sep"]

    # Map slide index to kind, dataset
    slide_map = {i: (k, ds) for (i, k, ds) in slides_info}

    # For all slides that are recognized image kinds
    for i, k, ds in slides_info:
        if k in by_kind:
            by_kind[k].append((ds if ds is not None else 10**9, i))

    # Sort each kind by dataset index (ascending)
    for k in by_kind:
        by_kind[k].sort(key=lambda x: x[0])

    # Compose new ordering of slide indices
    new_order = []
    # Keep the first title slide (if any) at the start
    if title_slide_indices:
        new_order.append(title_slide_indices[0])

    # Then group by type across subjects
    for k in type_order:
        # Optionally insert a header slide for the type (reuse the first encountered slide of that type as a template)
        # Simpler: just append the image slides directly in dataset order
        for _ds, idx in by_kind.get(k, []):
            new_order.append(idx)

    # Optionally move dataset separator slides to the end in original order
    new_order.extend(dataset_sep_indices)

    # Add any remaining slides not accounted for (unknown kinds)
    remaining = [i for (i, _k, _ds) in slides_info if i not in new_order]
    new_order.extend(remaining)

    # Reorder slides by manipulating underlying XML sldIdLst
    sldIdLst = prs.slides._sldIdLst
    sld_elems = list(sldIdLst)

    # Sanity check
    if len(sld_elems) != len(prs.slides):
        raise RuntimeError("Unexpected mismatch between slide XML and slides list")

    if sorted(new_order) != list(range(len(sld_elems))):
        raise RuntimeError("Computed new order does not reference all slides exactly once")

    # Clear and re-append in new order
    for _ in range(len(sld_elems)):
        sldIdLst.remove(sldIdLst[0])
    for idx in new_order:
        sldIdLst.append(sld_elems[idx])

    # Save
    if not output_path:
        base, ext = os.path.splitext(pptx_path)
        output_path = base + "_reordered_by_type" + ext
    prs.save(output_path)
    return output_path


def find_latest_combined_report(base_dir="/Users/cmilbourn/Documents/tSNR_check"):
    pats = glob(os.path.join(base_dir, "**/QA_Report_Combined_*.pptx"), recursive=True)
    if not pats:
        return None
    # Sort by timestamp in filename or modification time
    def extract_ts(p):
        m = re.search(r"QA_Report_Combined_(\d{8}_\d{6})\.pptx$", os.path.basename(p))
        if m:
            try:
                return datetime.strptime(m.group(1), "%Y%m%d_%H%M%S")
            except Exception:
                pass
        try:
            return datetime.fromtimestamp(os.path.getmtime(p))
        except Exception:
            return datetime.min
    pats.sort(key=extract_ts, reverse=True)
    return pats[0]


def main():
    if len(sys.argv) > 1:
        inp = sys.argv[1]
    else:
        inp = find_latest_combined_report()
        if not inp:
            print("Could not locate a QA_Report_Combined_*.pptx; please pass a path.")
            sys.exit(1)
    out = None
    if len(sys.argv) > 2:
        out = sys.argv[2]
    print(f"Reordering slides in: {inp}")
    output = reorder_by_type(inp, out)
    print(f"Wrote: {output}")

if __name__ == "__main__":
    main()
