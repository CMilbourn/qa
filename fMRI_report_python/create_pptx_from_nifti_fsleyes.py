#!/usr/bin/env python3
"""
Create PowerPoint from NIFTI files using fsleyes for visualization
Renders NIFTI data as images using fsleyes and compiles into PowerPoint
"""

import argparse
import os
import subprocess
from pathlib import Path
import json
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed.")

try:
    import nibabel as nib
    NIBABEL_AVAILABLE = True
except ImportError:
    NIBABEL_AVAILABLE = False
    print("⚠️  nibabel not installed.")

def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception:
        pass
    return None

def extract_metrics_from_nifti(nifti_path):
    """Extract metrics from NIFTI file"""
    metrics = {
        'tr': None,
        'shape': None,
        'affine': None,
        'data': None
    }
    
    # Get TR from JSON
    metrics['tr'] = get_tr_from_json(nifti_path)
    
    # Load NIFTI if nibabel available
    if NIBABEL_AVAILABLE:
        try:
            img = nib.load(nifti_path)
            metrics['shape'] = img.shape
            metrics['affine'] = img.affine
            # Load data for statistics
            data = img.get_fdata()
            metrics['data'] = data
        except Exception as e:
            print(f"⚠️  Could not load NIFTI {nifti_path}: {e}")
    
    return metrics

def render_nifti_with_fsleyes(nifti_path, output_png_path, scene_config=None):
    """
    Render a NIFTI file as an image using fsleyes
    
    Args:
        nifti_path: Path to the NIFTI file
        output_png_path: Path where to save the PNG
        scene_config: Optional fsleyes scene configuration
    
    Returns:
        True if successful, False otherwise
    """
    
    # Check if fsleyes is available
    try:
        result = subprocess.run(['which', 'fsleyes'], capture_output=True, text=True)
        if result.returncode != 0:
            print("⚠️  fsleyes not found in PATH")
            return False
    except Exception as e:
        print(f"⚠️  Error checking for fsleyes: {e}")
        return False
    
    try:
        # Basic fsleyes command to render NIFTI
        # Using renderScene to create a screenshot
        cmd = [
            'fsleyes',
            'render',
            nifti_path,
            '-of', output_png_path,
            '-s', '1600', '1200',  # Width and height as separate arguments
            '--hideCursor',
            '--hide', 'ui'
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
        
        if result.returncode == 0 and os.path.exists(output_png_path):
            return True
        else:
            print(f"⚠️  fsleyes render failed: {result.stderr}")
            return False
            
    except subprocess.TimeoutExpired:
        print(f"⚠️  fsleyes render timed out for {nifti_path}")
        return False
    except Exception as e:
        print(f"⚠️  Error rendering with fsleyes: {e}")
        return False

def create_simple_visualization(nifti_path, output_png_path):
    """
    Create a simple visualization using matplotlib as fallback if fsleyes unavailable
    Shows a middle slice of the volume
    """
    if not NIBABEL_AVAILABLE:
        print("⚠️  nibabel required for fallback visualization")
        return False
    
    try:
        import matplotlib.pyplot as plt
        
        img = nib.load(nifti_path)
        data = img.get_fdata()
        
        # Get middle slices
        if len(data.shape) == 4:
            # 4D data (fMRI)
            vol_idx = data.shape[3] // 2  # Middle timepoint
            vol_data = data[:, :, :, vol_idx]
        else:
            vol_data = data
        
        slice_idx = vol_data.shape[2] // 2
        slice_data = vol_data[:, :, slice_idx]
        
        fig, axes = plt.subplots(1, 3, figsize=(12, 4))
        
        # Axial
        axes[0].imshow(slice_data.T, cmap='gray', origin='lower')
        axes[0].set_title('Axial')
        axes[0].axis('off')
        
        # Sagittal
        slice_data_sag = vol_data[:, vol_data.shape[1] // 2, :]
        axes[1].imshow(slice_data_sag.T, cmap='gray', origin='lower')
        axes[1].set_title('Sagittal')
        axes[1].axis('off')
        
        # Coronal
        slice_data_cor = vol_data[vol_data.shape[0] // 2, :, :]
        axes[2].imshow(slice_data_cor.T, cmap='gray', origin='lower')
        axes[2].set_title('Coronal')
        axes[2].axis('off')
        
        plt.tight_layout()
        plt.savefig(output_png_path, dpi=100, bbox_inches='tight')
        plt.close()
        
        return os.path.exists(output_png_path)
        
    except Exception as e:
        print(f"⚠️  Error creating visualization: {e}")
        return False

def find_nifti_files(qa_dir, pattern='*nodummies.nii.gz'):
    """Find NIFTI files matching pattern in QA directory"""
    nifti_files = []
    try:
        for file in os.listdir(qa_dir):
            if file.endswith('.nii.gz') or file.endswith('.nii'):
                # Prefer nodummies files, but include others
                if 'nodummies' in file or len(nifti_files) == 0:
                    nifti_files.append(os.path.join(qa_dir, file))
    except Exception as e:
        print(f"⚠️  Error finding NIFTI files: {e}")
    
    return nifti_files

def add_metrics_table(slide, metrics_list, left=Inches(0.5), top=Inches(0.5)):
    """Add a table with metrics to the slide"""
    if not metrics_list:
        return
    
    rows = len(metrics_list) + 1
    cols = 3
    
    left_pos = left
    top_pos = top
    width = Inches(9)
    height = Inches(0.35) * rows
    
    table_shape = slide.shapes.add_table(rows, cols, left_pos, top_pos, width, height).table
    
    # Set column widths
    table_shape.columns[0].width = Inches(4)
    table_shape.columns[1].width = Inches(2.5)
    table_shape.columns[2].width = Inches(2.5)
    
    # Header row
    headers = ['Scan', 'TR (s)', 'Shape']
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 200, 200)
    
    # Data rows
    for row_idx, metrics in enumerate(metrics_list, start=1):
        # Scan name
        cell = table_shape.cell(row_idx, 0)
        cell.text = metrics.get('name', 'Unknown')
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        
        # TR
        tr_val = metrics.get('tr')
        cell = table_shape.cell(row_idx, 1)
        cell.text = f"{tr_val:.3f}s" if tr_val else "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Shape
        shape = metrics.get('shape')
        cell = table_shape.cell(row_idx, 2)
        if shape:
            shape_str = '×'.join(str(s) for s in shape)
            cell.text = shape_str
        else:
            cell.text = "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_powerpoint_from_nifti(input_dir, output_pptx, use_fsleyes=True):
    """
    Create PowerPoint from NIFTI files rendered as images
    """
    if not PPTX_AVAILABLE:
        print("❌ python-pptx is not installed")
        return False
    
    if not os.path.isdir(input_dir):
        print(f"❌ Input directory not found: {input_dir}")
        return False
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_images = 0
    all_metrics = []
    temp_pngs = []
    
    print("🔍 Scanning for QA directories...")
    
    # Iterate through subdirectories (QA output folders)
    subdirs = sorted([d for d in os.listdir(input_dir) 
                     if os.path.isdir(os.path.join(input_dir, d))])
    
    for subdir in subdirs:
        subdir_path = os.path.join(input_dir, subdir)
        
        print(f"\n📁 Processing: {subdir}")
        
        # Find NIFTI files
        nifti_files = find_nifti_files(subdir_path)
        
        if not nifti_files:
            print(f"   ⚠️  No NIFTI files found")
            continue
        
        # Process first NIFTI file
        nifti_path = nifti_files[0]
        print(f"   📊 Found NIFTI: {os.path.basename(nifti_path)}")
        
        # Extract metrics
        metrics = extract_metrics_from_nifti(nifti_path)
        metrics['name'] = subdir.split('_')[0]  # Extract scan ID
        all_metrics.append(metrics)
        
        # Create temporary PNG
        temp_png = os.path.join(subdir_path, f"_fsleyes_render_{os.path.basename(nifti_path)}.png")
        
        # Try fsleyes rendering
        success = False
        if use_fsleyes:
            print(f"   🎨 Rendering with fsleyes...")
            success = render_nifti_with_fsleyes(nifti_path, temp_png)
        
        # Fallback to matplotlib visualization
        if not success:
            print(f"   📈 Creating visualization with matplotlib...")
            success = create_simple_visualization(nifti_path, temp_png)
        
        if success:
            print(f"   ✅ Rendered: {os.path.basename(temp_png)}")
            temp_pngs.append((temp_png, subdir))
            total_images += 1
        else:
            print(f"   ❌ Failed to render")
    
    if not temp_pngs:
        print("\n❌ No NIFTI files were successfully rendered")
        return False
    
    # Create summary metrics slide
    if all_metrics:
        print("\n📊 Creating metrics summary slide...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "NIFTI Data Metrics Summary"
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        
        add_metrics_table(slide, all_metrics, left=Inches(0.3), top=Inches(0.8))
    
    # Create image slides
    print("\n🎨 Creating image slides...")
    for png_path, subdir_name in temp_pngs:
        try:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
            
            # Add title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = subdir_name.replace('_', ' ')
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            
            # Add image
            pic = slide.shapes.add_picture(png_path, Inches(0.5), Inches(1.0))
            
            # Scale image to fit
            max_width = Inches(9)
            max_height = Inches(6)
            
            if pic.width > max_width:
                ratio = max_width / pic.width
                pic.width = max_width
                pic.height = int(pic.height * ratio)
            
            if pic.height > max_height:
                ratio = max_height / pic.height
                pic.height = max_height
                pic.width = int(pic.width * ratio)
            
            print(f"   ✅ Added: {subdir_name}")
            
        except Exception as e:
            print(f"   ⚠️  Could not add image: {e}")
    
    # Save presentation
    try:
        prs.save(output_pptx)
        print(f"\n✅ PowerPoint saved: {output_pptx}")
        print(f"📊 Total images: {total_images}")
        
        # Clean up temp files
        print("\n🧹 Cleaning up temporary files...")
        for png_path, _ in temp_pngs:
            try:
                os.remove(png_path)
            except:
                pass
        
        return True
    except Exception as e:
        print(f"\n❌ Error saving PowerPoint: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(
        description='Create PowerPoint from NIFTI files using fsleyes visualization'
    )
    parser.add_argument(
        'input_dir', 
        help='Directory containing QA output subdirectories with NIFTI files'
    )
    parser.add_argument(
        '--output', 
        default='NIFTI_Report.pptx',
        help='Output PowerPoint filename (default: NIFTI_Report.pptx)'
    )
    parser.add_argument(
        '--no-fsleyes',
        action='store_true',
        help='Skip fsleyes and use matplotlib fallback only'
    )
    
    args = parser.parse_args()
    
    # Determine output path
    if os.path.isdir(args.input_dir):
        output_path = os.path.join(args.input_dir, args.output)
    else:
        output_path = args.output
    
    print(f"📁 Input directory: {args.input_dir}")
    print(f"📄 Output file: {output_path}")
    print(f"🎨 Using fsleyes: {not args.no_fsleyes}")
    print()
    
    success = create_powerpoint_from_nifti(args.input_dir, output_path, use_fsleyes=not args.no_fsleyes)
    
    if success:
        print("\n🎉 PowerPoint report created successfully!")
    else:
        print("\n❌ Failed to create PowerPoint")
    
    return 0 if success else 1

if __name__ == "__main__":
    exit(main())
