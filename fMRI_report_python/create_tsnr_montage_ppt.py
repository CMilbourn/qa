#!/usr/bin/env python3
"""
Create PowerPoint focusing specifically on tSNR Montage images (orange colormap)
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

def create_tsnr_focused_ppt(qa_parent_dir, output_file):
    """Create PowerPoint focused on tSNR montage images."""
    
    # Create presentation
    prs = Presentation()
    
    # Get all QA directories
    qa_dirs = []
    if os.path.exists(qa_parent_dir):
        for item in os.listdir(qa_parent_dir):
            item_path = os.path.join(qa_parent_dir, item)
            if os.path.isdir(item_path) and item.startswith('qa_output'):
                qa_dirs.append((item, item_path))
    
    qa_dirs.sort()  # Sort alphabetically
    
    if not qa_dirs:
        print(f"❌ No QA directories found in {qa_parent_dir}")
        return False
    
    print(f"📊 Creating tSNR-focused PowerPoint with {len(qa_dirs)} datasets")
    
    # Title slide
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Temporal SNR (tSNR) Analysis Report"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"sub003 tSNR Montage Analysis - {len(qa_dirs)} Acquisitions"
    
    # Summary slide with tSNR metrics
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "tSNR Summary Overview"
    
    summary_text = "Temporal SNR Analysis Summary:\\n\\n"
    
    for i, (qa_name, qa_dir) in enumerate(qa_dirs, 1):
        # Extract scan info from name
        if "2mm_longerTR" in qa_name:
            scan_type = "2mm Resolution (Long TR)"
        elif "3mm" in qa_name:
            scan_type = "3mm Resolution"
        elif "2mm_pre" in qa_name and "de-9" in qa_name:
            scan_type = "2mm Resolution (Pre-scan 1)"
        elif "2mm_pre" in qa_name and "de-2" in qa_name:
            scan_type = "2mm Resolution (Pre-scan 2)"
        else:
            scan_type = "Unknown Resolution"
        
        summary_text += f"{i}. {scan_type}\\n"
        
        # Look for tSNR files and metrics
        tsnr_files = []
        for root, dirs, files in os.walk(qa_dir):
            for file in files:
                if 'tsnr' in file.lower() or 'tSNR' in file:
                    if file.endswith('.png'):
                        tsnr_files.append(file)
        
        summary_text += f"   • tSNR images: {len(tsnr_files)}\\n"
        summary_text += f"   • Directory: {os.path.basename(qa_dir)[:50]}...\\n\\n"
    
    slide.placeholders[1].text = summary_text
    
    # Create detailed slides for each dataset
    tsnr_image_types = [
        ('tSNR_montage.png', 'tSNR Montage (Orange Colormap)'),
        ('tSNR_w_ROI_images.png', 'tSNR with ROI Analysis'),
        ('tSNR_per_unit_time.png', 'tSNR per Unit Time'),
        ('tSNR_raw.png', 'Raw tSNR'),
        ('tSNR_cor.png', 'tSNR Coronal View'),
        ('tSNR_sag.png', 'tSNR Sagittal View')
    ]
    
    for qa_name, qa_dir in qa_dirs:
        print(f"Processing: {qa_name}")
        
        # Section title slide for each acquisition
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        if "2mm_longerTR" in qa_name:
            scan_title = "2mm Resolution (Long TR) - tSNR Analysis"
        elif "3mm" in qa_name:
            scan_title = "3mm Resolution - tSNR Analysis"
        elif "2mm_pre" in qa_name and "de-9" in qa_name:
            scan_title = "2mm Resolution (Pre-scan 1) - tSNR Analysis"
        elif "2mm_pre" in qa_name and "de-2" in qa_name:
            scan_title = "2mm Resolution (Pre-scan 2) - tSNR Analysis"
        else:
            scan_title = "fMRI Acquisition - tSNR Analysis"
        
        slide.shapes.title.text = scan_title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = f"Temporal Signal-to-Noise Ratio Visualization\\nOrange colormap indicates higher tSNR values"
        
        # Add tSNR image slides
        for image_file, image_title in tsnr_image_types:
            image_path = os.path.join(qa_dir, image_file)
            
            if os.path.exists(image_path):
                # Use blank layout for image slides
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title at top
                title_left = Inches(0.5)
                title_top = Inches(0.2)
                title_width = Inches(9)
                title_height = Inches(0.8)
                
                title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_frame = title_shape.text_frame
                title_frame.text = f"{scan_title}: {image_title}"
                
                # Format title
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(20)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Add image (preserve aspect ratio, centered)
                try:
                    img_left = Inches(1)
                    img_top = Inches(1.2)
                    img_width = Inches(8)
                    
                    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
                    
                except Exception as e:
                    print(f"  ⚠️  Error adding image {image_file}: {e}")
                
                # Add description for tSNR montage
                if 'montage' in image_file.lower():
                    desc_left = Inches(0.5)
                    desc_top = Inches(6.5)
                    desc_width = Inches(9)
                    desc_height = Inches(1)
                    
                    desc_shape = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
                    desc_frame = desc_shape.text_frame
                    desc_frame.text = "Orange/Red colors indicate higher temporal SNR values. Brighter regions show better signal stability over time."
                    
                    # Format description text
                    desc_paragraph = desc_frame.paragraphs[0]
                    desc_paragraph.font.size = Pt(12)
                    desc_paragraph.alignment = PP_ALIGN.CENTER
                
                print(f"  ✅ Added: {image_title}")
            else:
                print(f"  ⚠️  Missing: {image_file}")
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✅ tSNR-focused PowerPoint saved: {output_file}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        print(f"📊 File size: {file_size:.1f} MB")
        print(f"📄 Total slides: {len(prs.slides)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    """Main function."""
    
    qa_parent_dir = "/Users/cmilbourn/Documents/Sweet_Data/Development_Data/sub003/sub003-visit001-ses001/func/"
    output_file = "/Users/cmilbourn/Documents/Sweet_Data/Development_Data/sub003/sub003-visit001-ses001/func/20251015_sub003_tSNR_MONTAGE_report.pptx"
    
    print(f"🚀 Creating tSNR-focused PowerPoint report from: {qa_parent_dir}")
    
    success = create_tsnr_focused_ppt(qa_parent_dir, output_file)
    
    if success:
        print(f"🎉 tSNR montage report generation complete!")
        print(f"📁 Saved as: {os.path.basename(output_file)}")
    else:
        print(f"❌ Report generation failed!")

if __name__ == "__main__":
    main()