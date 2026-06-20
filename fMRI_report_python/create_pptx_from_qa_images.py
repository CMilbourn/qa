#!/usr/bin/env python3
"""
Create PowerPoint from QA output images with metric tables
Compiles PNG images from QA directories into a PowerPoint presentation
Includes tables with TR, tSNR, and tSNR per unit time metrics
"""

import argparse
import os
import json
from pathlib import Path
import numpy as np

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️  python-pptx not installed.")

try:
    import nibabel as nib
    NIBABEL_AVAILABLE = True
except ImportError:
    NIBABEL_AVAILABLE = False

def get_tr_from_json(nifti_path):
    """Extract TR from corresponding JSON file"""
    try:
        json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
        if os.path.exists(json_path):
            with open(json_path, 'r') as f:
                metadata = json.load(f)
                tr = metadata.get('RepetitionTime')
                if tr:
                    return float(tr)
    except Exception as e:
        pass
    return None

def extract_metrics_from_qa_dir(qa_dir):
    """Extract TR and tSNR metrics from a QA output directory"""
    metrics = {
        'tr': None,
        'tsnr': None,
        'tsnr_per_unit_time': None,
        'tsnr_masked': None
    }
    
    # Try to find NIFTI file and extract TR from JSON
    for file in os.listdir(qa_dir):
        if file.endswith(('.nii.gz', '.nii')):
            tr = get_tr_from_json(os.path.join(qa_dir, file))
            if tr:
                metrics['tr'] = tr
            break
    
    # Try to extract tSNR from NIFTI files
    if NIBABEL_AVAILABLE:
        tsnr_file = os.path.join(qa_dir, 'tsnr.nii.gz')
        if os.path.exists(tsnr_file):
            try:
                tsnr_img = nib.load(tsnr_file)
                tsnr_data = tsnr_img.get_fdata()
                tsnr_value = float(np.mean(tsnr_data[tsnr_data > 0]))
                metrics['tsnr'] = tsnr_value
                
                if metrics['tr'] and tsnr_value:
                    metrics['tsnr_per_unit_time'] = tsnr_value / np.sqrt(metrics['tr'])
            except Exception:
                pass
    
    return metrics

def get_png_files_from_dir(directory):
    """Get all PNG files from a directory, sorted"""
    png_files = []
    if os.path.exists(directory):
        for filename in sorted(os.listdir(directory)):
            if filename.endswith('.png'):
                png_files.append(os.path.join(directory, filename))
    return png_files

def add_metrics_table(slide, metrics_list, left=Inches(0.5), top=Inches(0.5)):
    """Add a table with metrics to the slide"""
    rows = len(metrics_list) + 1
    cols = 4
    
    left_pos = left
    top_pos = top
    width = Inches(9)
    height = Inches(0.4) * rows
    
    table_shape = slide.shapes.add_table(rows, cols, left_pos, top_pos, width, height).table
    
    # Set column widths
    table_shape.columns[0].width = Inches(3.5)
    table_shape.columns[1].width = Inches(1.5)
    table_shape.columns[2].width = Inches(2)
    table_shape.columns[3].width = Inches(2)
    
    # Header row
    headers = ['Scan', 'TR (s)', 'tSNR', 'tSNR/√TR']
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        # Format header
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(9)
        paragraph.alignment = PP_ALIGN.CENTER
        # Header background color
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200, 200, 200)
    
    # Data rows
    for row_idx, metrics in enumerate(metrics_list, start=1):
        # Scan name
        cell = table_shape.cell(row_idx, 0)
        cell.text = metrics.get('name', 'Unknown')
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        
        # TR
        tr_val = metrics.get('tr')
        cell = table_shape.cell(row_idx, 1)
        cell.text = f"{tr_val:.3f}" if tr_val else "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # tSNR
        tsnr_val = metrics.get('tsnr')
        cell = table_shape.cell(row_idx, 2)
        cell.text = f"{tsnr_val:.2f}" if tsnr_val else "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # tSNR per unit time
        tsnr_unit = metrics.get('tsnr_per_unit_time')
        cell = table_shape.cell(row_idx, 3)
        cell.text = f"{tsnr_unit:.2f}" if tsnr_unit else "N/A"
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_powerpoint_from_images(input_dir, output_pptx):
    """
    Create a PowerPoint presentation from PNG images in a directory
    Images are organized by subdirectory with metric tables
    """
    if not PPTX_AVAILABLE:
        print("❌ python-pptx is not installed")
        return False
    
    if not os.path.isdir(input_dir):
        print(f"❌ Input directory not found: {input_dir}")
        return False
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    total_images = 0
    all_metrics = []
    
    # First pass: collect metrics
    subdirs = sorted([d for d in os.listdir(input_dir) if os.path.isdir(os.path.join(input_dir, d))])
    
    for subdir in subdirs:
        subdir_path = os.path.join(input_dir, subdir)
        metrics = extract_metrics_from_qa_dir(subdir_path)
        metrics['name'] = subdir.split('_')[0]  # Extract just the scan ID
        all_metrics.append(metrics)
    
    # Create summary metrics slide
    if all_metrics:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = "QA Metrics Summary"
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        
        add_metrics_table(slide, all_metrics, left=Inches(0.3), top=Inches(0.8))
    
    # Second pass: create image slides
    for subdir in subdirs:
        subdir_path = os.path.join(input_dir, subdir)
        
        # Get PNG files from this subdirectory
        png_files = get_png_files_from_dir(subdir_path)
        
        if not png_files:
            continue
        
        print(f"📁 Processing: {subdir}")
        print(f"   Found {len(png_files)} images")
        
        # Create a new slide for this subdirectory
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = subdir.replace('_', ' ')
        title_para.font.size = Pt(20)
        title_para.font.bold = True
        
        # Add images to the slide (up to 4 per slide)
        img_y_position = Inches(1.0)
        images_on_slide = 0
        
        for png_file in png_files:
            try:
                # Check if we need a new slide
                if images_on_slide >= 4:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    # Add subtitle indicating it's continuation
                    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.4))
                    subtitle_frame = subtitle_box.text_frame
                    subtitle_para = subtitle_frame.paragraphs[0]
                    subtitle_para.text = f"{subdir.replace('_', ' ')} (continued)"
                    subtitle_para.font.size = Pt(14)
                    img_y_position = Inches(1.0)
                    images_on_slide = 0
                
                # Add image
                pic = slide.shapes.add_picture(png_file, Inches(0.5), img_y_position)
                
                # Scale image to fit
                max_width = Inches(9)
                max_height = Inches(1.4)
                
                if pic.width > max_width:
                    ratio = max_width / pic.width
                    pic.width = max_width
                    pic.height = int(pic.height * ratio)
                
                if pic.height > max_height:
                    ratio = max_height / pic.height
                    pic.height = max_height
                    pic.width = int(pic.width * ratio)
                
                # Add caption
                caption_text = os.path.basename(png_file).replace('.png', '').replace('_', ' ')
                caption_box = slide.shapes.add_textbox(
                    Inches(0.5), 
                    img_y_position + pic.height + Inches(0.05), 
                    Inches(9), 
                    Inches(0.25)
                )
                caption_frame = caption_box.text_frame
                caption_para = caption_frame.paragraphs[0]
                caption_para.text = caption_text
                caption_para.font.size = Pt(8)
                
                img_y_position = img_y_position + pic.height + Inches(0.4)
                images_on_slide += 1
                total_images += 1
                
            except Exception as e:
                print(f"⚠️  Could not add image {png_file}: {e}")
    
    if total_images == 0:
        print("❌ No PNG images found")
        return False
    
    # Save presentation
    try:
        prs.save(output_pptx)
        print(f"✅ PowerPoint saved: {output_pptx}")
        print(f"📊 Total images: {total_images}")
        return True
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    parser = argparse.ArgumentParser(
        description='Create PowerPoint from QA output images with metrics'
    )
    parser.add_argument(
        'input_dir', 
        help='Directory containing QA output subdirectories with PNG images'
    )
    parser.add_argument(
        '--output', 
        default='QA_Report.pptx',
        help='Output PowerPoint filename (default: QA_Report.pptx)'
    )
    
    args = parser.parse_args()
    
    # Determine output path
    if os.path.isdir(args.input_dir):
        output_path = os.path.join(args.input_dir, args.output)
    else:
        output_path = args.output
    
    print(f"📁 Input directory: {args.input_dir}")
    print(f"📄 Output file: {output_path}")
    print()
    
    success = create_powerpoint_from_images(args.input_dir, output_path)
    
    if success:
        print("\n🎉 PowerPoint report created successfully!")
    else:
        print("\n❌ Failed to create PowerPoint")
    
    return 0 if success else 1

if __name__ == "__main__":
    exit(main())
