#!/usr/bin/env python3
"""
Simple PowerPoint Creator for sub001 QA Results

Creates a comprehensive PowerPoint presentation from all QA output directories.
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import json

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_content_slide(prs, title, content=""):
    """Add a content slide with title."""
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if content:
        slide.placeholders[1].text = content
    
    return slide

def add_image_slide(prs, title, image_path, description=""):
    """Add a slide with an image."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(24)
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Add image (preserve aspect ratio)
    if os.path.exists(image_path):
        try:
            # Calculate position to center the image
            img_left = Inches(1)
            img_top = Inches(1.2)
            img_width = Inches(8)
            
            slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
            
        except Exception as e:
            print(f"Error adding image {image_path}: {e}")
    
    # Add description if provided
    if description:
        desc_shape = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
        desc_frame = desc_shape.text_frame
        desc_frame.text = description
        desc_paragraph = desc_frame.paragraphs[0]
        desc_paragraph.font.size = Pt(12)
        desc_paragraph.alignment = PP_ALIGN.CENTER

def extract_qa_summary(qa_dir):
    """Extract QA summary information from directory."""
    summary_file = os.path.join(qa_dir, 'qa_summary.json')
    
    if os.path.exists(summary_file):
        with open(summary_file, 'r') as f:
            return json.load(f)
    else:
        return {}

def create_comprehensive_ppt(qa_parent_dir, output_file):
    """Create comprehensive PowerPoint from all QA directories."""
    
    # Create presentation
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, "fMRI Quality Assessment Report", "sub001 Complete Analysis")
    
    # Get all QA directories
    qa_dirs = []
    if os.path.exists(qa_parent_dir):
        for item in os.listdir(qa_parent_dir):
            item_path = os.path.join(qa_parent_dir, item)
            if os.path.isdir(item_path):
                qa_dirs.append((item, item_path))
    
    qa_dirs.sort()  # Sort alphabetically
    
    if not qa_dirs:
        print(f"❌ No QA directories found in {qa_parent_dir}")
        return False
    
    print(f"📊 Found {len(qa_dirs)} QA result directories")
    
    # Summary slide
    summary_slide = add_content_slide(prs, "Processing Summary", "")
    
    summary_text = f"Total scans processed: {len(qa_dirs)}\\n\\n"
    
    for qa_name, qa_dir in qa_dirs:
        summary_info = extract_qa_summary(qa_dir)
        
        # Extract key info from filename
        if "2mm" in qa_name and "1.5" not in qa_name:
            scan_type = "2mm Resolution"
        elif "2mm_1.5" in qa_name:
            scan_type = "2mm x 1.5mm Resolution"
        elif "1.5_mm_iso" in qa_name:
            scan_type = "1.5mm Isotropic Resolution"
        else:
            scan_type = "Unknown Resolution"
        
        tr = summary_info.get('tr', 'Unknown')
        mean_tsnr = summary_info.get('mean_tsnr', 0)
        mean_isnr = summary_info.get('mean_isnr', 0)
        
        summary_text += f"• {scan_type}\\n"
        summary_text += f"  TR: {tr}s, tSNR: {mean_tsnr:.1f}, iSNR: {mean_isnr:.0f}\\n\\n"
    
    summary_slide.placeholders[1].text = summary_text
    
    # Add slides for each QA directory
    image_types = [
        ('mean_montage.png', 'Mean Signal Montage'),
        ('tsnr_montage.png', 'Temporal SNR Montage'),
        ('isnr_montage.png', 'Instantaneous SNR Montage'),
        ('sample_timeseries.png', 'Sample Time Series')
    ]
    
    for qa_name, qa_dir in qa_dirs:
        print(f"Processing: {qa_name}")
        
        # Extract scan info
        summary_info = extract_qa_summary(qa_dir)
        
        # Section title slide
        if "2mm" in qa_name and "1.5" not in qa_name:
            scan_title = "2mm Resolution Scan"
        elif "2mm_1.5" in qa_name:
            scan_title = "2mm x 1.5mm Resolution Scan"
        elif "1.5_mm_iso" in qa_name:
            scan_title = "1.5mm Isotropic Resolution Scan"
        else:
            scan_title = "fMRI Scan"
        
        tr = summary_info.get('tr', 'Unknown')
        shape = summary_info.get('shape', [0,0,0,0])
        ernst_scaling = summary_info.get('ernst_scaling', 1.0)
        
        subtitle = f"TR: {tr}s | Matrix: {shape[0]}×{shape[1]}×{shape[2]}×{shape[3]} | Ernst Scaling: {ernst_scaling:.3f}"
        add_title_slide(prs, scan_title, subtitle)
        
        # Add image slides
        for image_file, image_title in image_types:
            image_path = os.path.join(qa_dir, image_file)
            
            if os.path.exists(image_path):
                full_title = f"{scan_title}: {image_title}"
                
                # Add metrics to description for SNR images
                description = ""
                if 'tsnr' in image_file.lower():
                    tsnr = summary_info.get('mean_tsnr', 0)
                    description = f"Mean tSNR: {tsnr:.2f}"
                elif 'isnr' in image_file.lower():
                    isnr = summary_info.get('mean_isnr', 0)
                    description = f"Mean iSNR: {isnr:.0f}"
                
                add_image_slide(prs, full_title, image_path, description)
                print(f"  Added: {image_title}")
            else:
                print(f"  ⚠️  Missing: {image_file}")
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✅ PowerPoint saved: {output_file}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        print(f"📊 File size: {file_size:.1f} MB")
        print(f"📄 Total slides: {len(prs.slides)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    """Main function."""
    if len(sys.argv) < 3:
        print("Usage: python create_sub001_ppt.py <qa_batch_dir> <output_file>")
        print("Example: python create_sub001_ppt.py qa_output_batch sub001_QA_report.pptx")
        sys.exit(1)
    
    qa_parent_dir = sys.argv[1]
    output_file = sys.argv[2]
    
    if not os.path.exists(qa_parent_dir):
        print(f"❌ QA directory not found: {qa_parent_dir}")
        sys.exit(1)
    
    print(f"🚀 Creating PowerPoint report from: {qa_parent_dir}")
    
    success = create_comprehensive_ppt(qa_parent_dir, output_file)
    
    if success:
        print(f"🎉 Report generation complete!")
    else:
        print(f"❌ Report generation failed!")
        sys.exit(1)

if __name__ == "__main__":
    main()