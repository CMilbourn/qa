#!/usr/bin/env python
# coding: utf-8
# qa_run_nophase.py based on qa_run_nophase.ipynb  
# example run line: python qa_run_nophase_V9_multisubj_MASK_30dyn.py
# 20251218 add the run for running the tsnr round what ever it is now with Sue F
#reorder_qa_pptx_by_type.py
# # Quality Assurance (QA) Python Version
# 
# This notebook gives example uses of image based QA metrics in `qa`.
# 
# A large chunk of this code has been taken from Alex Daniel's `ukat` code: https://github.com/UKRIN-MAPS/ukat.
# 
# We'll start with some imports and general housekeeping.
# 
# Reference for tSNR: https://doi.org/10.1016/j.neuroimage.2005.01.007

# ## READ THIS INFO
# 
# - Scroll down to the bottom to set paths etc. The key things to change are
# `mypathname = '/Users/spmic/data/tw_testing_may_2025/'`
# `pathname_m = mypathname + 'magnitude/'`
# `extension = '.nii' #this can be .nii or .nii.gz`
# `filename_pattern = 'digitmap*'`
# 
# - Just add your path to the `qa` github folder here
# 
# - There is a line `noise_volume = imgm_cla[:, :, :, -1]` This grabs the last scan, assumes it is noise and uses it for iSNR.
# 
# - In the function `process_data_nophase` there are options to set slices/scales/sizes for patch ROI etc. near the top
# 
# - There is a bit that calculates the tSNR per unit time using the TR, this makes a separate PNG, can ignore this or set the TR yourself. Shouldn't affect anything else.
# 
# - If you have a nifti with the word `mask` in it, in the same folder, then it will find this and mask your data by it.

import sys
sys.path.append('/Users/cmilbourn/Documents/GitHub/qa/')  # ** change line to match code folder location **
print(sys.path)

import os
import numpy as np
import matplotlib
matplotlib.use('Agg')  # Use non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.patches as patches
from matplotlib.colors import Normalize
import nibabel as nib
from glob import glob
import subprocess
import json
import csv
from datetime import datetime
import argparse
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
    PPTX_AVAILABLE = True
except ImportError:
    print('python-pptx not installed, PowerPoint generation will be skipped.')
    PPTX_AVAILABLE = False

#from ukat.data import fetch
from fMRI_report_python.functions import snr
#from fMRI_report_python.functions.snr import some_function
#from functions import snr
from scipy.signal import detrend

from mpl_toolkits.mplot3d import Axes3D  # Needed for 3D plotting

# check packages
print(snr)
print(dir(snr))

# make the loading a bit easier inline
def load_data(inputdatafilename):
    #Function to load in data 
    #Needs an input arg
    data = nib.load(inputdatafilename)
    image = data.get_fdata()
    return image, data.affine

# Function to find a mask file for a given dataset
def find_mask_file(directory, mag_basename=None):
    """
    Prefer a per-file mask in the same directory that matches the magnitude filename.
    Fallbacks: any *_brain_mask.nii[.gz], then any *mask.nii[.gz] in the directory.

    Parameters
    ----------
    directory : str
        Directory to search for masks
    mag_basename : str or None
        Basename (without extension) of the magnitude file to match, e.g.,
        'foo_30dyn' so we can prefer 'foo_30dyn_brain_mask.nii.gz'.
    """
    try:
        candidates = os.listdir(directory)
    except Exception:
        return None

    def _exists(name):
        return os.path.exists(os.path.join(directory, name))

    # 1) Exact expected names based on mag_basename
    if mag_basename:
        expected = []
        # If input already ends with _30dyn, prefer that exact suffix
        expected.append(f"{mag_basename}_brain_mask.nii.gz")
        expected.append(f"{mag_basename}_brain_mask.nii")
        # Some BET outputs may put _30dyn before suffix; we already included mag_basename
        for fname in expected:
            if fname in candidates and _exists(fname):
                return os.path.join(directory, fname)

    # 2) Any brain_mask in directory
    for filename in candidates:
        fl = filename.lower()
        if "brain_mask" in fl and fl.endswith((".nii", ".nii.gz")):
            return os.path.join(directory, filename)

    # 3) Any mask in directory
    for filename in candidates:
        fl = filename.lower()
        if "mask" in fl and fl.endswith((".nii", ".nii.gz")):
            return os.path.join(directory, filename)

    return None

# Function to read TR from BIDS JSON sidecar file
def get_tr_from_json(nifti_path):
    """
    Read RepetitionTime from BIDS JSON sidecar file
    
    Parameters:
    -----------
    nifti_path : str
        Path to the NIfTI file (e.g., sub001-task-rest-bold.nii.gz)
        
    Returns:
    --------
    float or None
        RepetitionTime in seconds, or None if not found
    """
    # Construct JSON filename by replacing .nii.gz or .nii with .json
    json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
    
    if os.path.exists(json_path):
        try:
            with open(json_path, 'r') as f:
                metadata = json.load(f)
            
            tr = metadata.get('RepetitionTime')
            if tr is not None:
                print(f"Found TR = {tr}s in JSON file: {json_path}")
                return float(tr)
            else:
                print(f"RepetitionTime not found in JSON file: {json_path}")
                return None
                
        except Exception as e:
            print(f"Error reading JSON file {json_path}: {e}")
            return None
    else:
        print(f"No JSON sidecar found at: {json_path}")
        return None

# Function to read TE from BIDS JSON sidecar file
def get_te_from_json(nifti_path):
    """
    Read EchoTime from BIDS JSON sidecar file
    
    Parameters:
    -----------
    nifti_path : str
        Path to the NIfTI file
        
    Returns:
    --------
    float or None
        EchoTime in seconds, or None if not found
    """
    json_path = nifti_path.replace('.nii.gz', '.json').replace('.nii', '.json')
    
    if os.path.exists(json_path):
        try:
            with open(json_path, 'r') as f:
                metadata = json.load(f)
            
            te = metadata.get('EchoTime')
            if te is not None:
                print(f"Found TE = {te}s in JSON file: {json_path}")
                return float(te)
            else:
                print(f"EchoTime not found in JSON file: {json_path}")
                return None
                
        except Exception as e:
            print(f"Error reading JSON file {json_path}: {e}")
            return None
    else:
        return None

# Function to get voxel size from NIfTI header
def get_voxel_size(nifti_img):
    """
    Extract voxel dimensions from NIfTI image
    
    Parameters:
    -----------
    nifti_img : nibabel Nifti1Image
        Loaded NIfTI image object
        
    Returns:
    --------
    tuple or None
        (x, y, z) voxel dimensions in mm, or None if not found
    """
    try:
        # Get voxel dimensions from header (pixdim)
        header = nifti_img.header
        zooms = header.get_zooms()
        # Return first 3 dimensions (x, y, z spatial dimensions)
        return tuple(zooms[:3])
    except Exception as e:
        print(f"Error getting voxel size: {e}")
        return None

def process_data_nophase(imgm_cla, imgm_affine, core_filename, output_dir, mask_data=None, TR=None, nifti_path=None, nifti_img=None):
    # This is our big function
    # I had to condense all the functionality of the notebook into this one func, to make the for loop in the script easier
    # So, first it plots the mean images
    # Calculates iSNR
    # Calculates tSNR
    # Calculates tSNR in a patch ROI
    # Plots signal and std of signal in patch over time
    # Plots static spatial noise image
    # OUTPUTS: 
    # tSNR images saved as nii.gz
    # Lots of PNG images of every plot

    ############################## Plotting mean images
    # Set slice (3d) and time (4d) - use middle slices for all dimensions
    slice_index = imgm_cla.shape[2] // 2  # Middle axial slice
    slice_index_sag = imgm_cla.shape[0] // 2  # Middle sagittal slice
    slice_index_cor = imgm_cla.shape[1] // 2  # Middle coronal slice
    print(f"Slice indices - Axial: {slice_index}, Sagittal: {slice_index_sag}, Coronal: {slice_index_cor}")
    print(f"Image dimensions: {imgm_cla.shape}")
    time_point = 1
    tsnrScale = 100
    #tsnrScale = 10
    #isnrScale = 100  
    isnrScale = 10

    # We want to define a patch ROI
    # Define parameters for ROI size and position
    x_start = 60  # Starting x-coordinate of the ROI
    y_start = 80 # Starting y-coordinate of the ROI
    roi_width = 20  # Width of the ROI
    roi_height = 20  # Height of the ROI

    # if mask_data is not None:
    #     masked_img = imgm_cla[mask_data > 0]

    mean_img = np.mean(imgm_cla, axis=3)
    
    # Create a 2x2 grid of subplots
    fig, axs = plt.subplots(1, 1, figsize=(5, 5))  # Adjust figsize as needed

    # Plot imgm_cla
    axs.imshow(mean_img[:, :, slice_index].T, origin='lower', cmap='gray')
    axs.set_title(f'Magnitude (Slice {slice_index})')  # Set subplot title with slice index
    axs.axis(False)  # Turn off axis labels and ticks

    # Adjust layout and display the plot
    plt.tight_layout()
    # # plt.show()  # Disabled for script mode  # Disabled for script mode

    # Assuming you have created a plot `plt` and want to save it as a PNG file
    output_filename = 'Mean_image.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    # Determine grid size for montage
    num_slices = imgm_cla.shape[2]
    rows = int(np.ceil(np.sqrt(num_slices)))
    cols = int(np.ceil(num_slices / rows))
    
    # Create figure with larger size for bigger individual plots
    figsize_scale = 2.5  # Scale this to make plots larger
    fig = plt.figure(figsize=(cols * figsize_scale, rows * figsize_scale))
    
    # Calculate contrast limits based on percentiles to improve visibility
    vmin = np.percentile(mean_img, 2)
    vmax = np.percentile(mean_img, 98)
    
    # Loop through slices
    im_list = []
    for i in range(num_slices):
        ax = fig.add_subplot(rows, cols, i + 1)
        im = ax.imshow(mean_img[:, :, i], cmap='gray', vmin=vmin, vmax=vmax)
        if i == 0:
            im_list.append(im)
        ax.set_title(f"Slice {i}", fontsize=8)  # Smaller font size
        ax.axis('off')
    
    # Add colorbar on the right side
    fig.subplots_adjust(right=0.90)
    cbar_ax = fig.add_axes([0.92, 0.15, 0.02, 0.7])
    fig.colorbar(im_list[0], cax=cbar_ax, label='Intensity')
    
    # Tighter layout
    fig.tight_layout(pad=0.3, rect=[0, 0, 0.90, 1])
    
    # Save output
    output_filename = 'mean_montage.png'
    output_path = f"{output_dir}/{output_filename}"
    plt.savefig(output_path, dpi=300, bbox_inches='tight')

    ############################## iSNR
    # ISNR show    

    # Compute iSNR
    # isnr_test = snr.Isnr(imgm_cla, imgm_cla_affine)
    # # Extract only the first volume from the iSNR map while keeping 4D structure
    # isnr_test.isnr_map = isnr_test.isnr_map[:, :, :, 0:1]  # Keeps it 4D
    # # Save as NIFTI
    # isnr_test.to_nifti(OUTPUT_DIR, 'isnr')
    # print(f'Image has an iSNR of {isnr_test.isnr:.2f}.', file=output_file)

    #isnr_cla = snr.Isnr(imgm_cla, imgm_cla_affine).isnr
    #isnr_test = snr.Isnr(imgm_cla, imgm_cla_affine)
    #isnr_test = isnr_test[:, :, :, 0]
    # Save as NIFTI 
    #isnr_test.to_nifti(OUTPUT_DIR, 'isnr')
    #print(f'Image has an iSNR of {isnr_cla:.2f}.', file=output_file)

    noise_volume = imgm_cla[:, :, :, -1]  # Extract the last volume
    # if no noise scan, use difference of 2 dynamics
    # Assuming the volumes you want to subtract are at indices 15 and 16
    #vol1 = imgm_cla[:, :, :, 3]  # Extract volume 15
    #vol2 = imgm_cla[:, :, :, 4]  # Extract volume 16
    # # Calculate the difference between the two volumes
    #noise_volume = vol1 - vol2
    # # Replace NaNs with 0
    noise_volume = np.nan_to_num(noise_volume, nan=0.0)

    
    # Compute the temporal mean image
    tmean_img = np.mean(imgm_cla, axis=-1)
    # Create a brain mask by thresholding (adjust threshold as needed)
    brain_mask = tmean_img > (0.05 * np.max(tmean_img))  # Example: 10% of max intensity
    
    # Apply the mask to the noise volume
    masked_noise = noise_volume * brain_mask

    masked_mean = mean_img * brain_mask
    valid_voxels = masked_mean[masked_mean != 0]  # Exclude zeros
    mean_std = np.std(valid_voxels)  # Compute stddev only for non-zero voxels
    print(f'Mean volume std = {mean_std:.2f}.')
    
    fig, axs = plt.subplots(1, 3, figsize=(15, 5))  # 3 subplots in one row
    # Plot Noise Volume
    im1 = axs[0].imshow(noise_volume[:, :, slice_index].T, origin='lower', cmap='viridis')
    axs[0].set_title("Noise Volume")
    axs[0].axis(False)
    fig.colorbar(im1, ax=axs[0])
    
    # Plot Brain Mask
    im2 = axs[1].imshow(brain_mask[:, :, slice_index].T, origin='lower', cmap='gray')
    axs[1].set_title("Brain Mask")
    axs[1].axis(False)
    fig.colorbar(im2, ax=axs[1])
    
    # Plot Masked Noise
    im3 = axs[2].imshow(masked_noise[:, :, slice_index].T, origin='lower', cmap='viridis')
    axs[2].set_title("Masked Noise")
    axs[2].axis(False)
    fig.colorbar(im3, ax=axs[2])
    # Adjust layout and display
    plt.tight_layout()
    # plt.show()  # Disabled for script mode
    # Save the image
    output_filename = 'masked_noise.png'
    output_path = f"{output_dir}/{output_filename}"
    fig.savefig(output_path, dpi=300)  
    plt.close()  # Close to free up memory
    
    
    isnr_obj_cla = snr.Isnr(imgm_cla, imgm_affine, noise_mask=masked_noise)
    isnr_obj_cla.to_nifti(output_dir, 'isnr')
    print(f'Image has an iSNR of {np.mean(isnr_obj_cla.isnr):.2f}.')
    print(f'Noise value used for iSNR: {np.mean(isnr_obj_cla.noise):.2f}')

    # Plot Noise volume montage
    # Define the number of rows and columns for montage
    num_slices = noise_volume.shape[2]  # Number of slices
    
    # Create a figure for the montage
    fig = plt.figure(figsize=(15, 15))  # Adjust figure size
    im_noise = None
    for i in range(num_slices):
        ax = fig.add_subplot(rows, cols, i + 1)  # Create subplot for each slice
        im = ax.imshow(noise_volume[:, :, i], cmap='gray', clim=(0, 2000))  # Display slice
        if i == 0:
            im_noise = im
        ax.set_title(f"Slice {i}")  # Label each slice
        ax.axis('off')  # Hide axes
    # Add colorbar on the right side
    fig.subplots_adjust(right=0.90)
    cbar_ax = fig.add_axes([0.92, 0.15, 0.02, 0.7])
    fig.colorbar(im_noise, cax=cbar_ax, label='Intensity')
    # Adjust layout
    fig.tight_layout(pad=0.5, rect=[0, 0, 0.90, 1])
    ## plt.show()  # Disabled for script mode
    # Save the montage
    output_filename = 'noise_volume_montage.png'
    output_path = f"{output_dir}/{output_filename}"
    plt.savefig(output_path, dpi=300)
    #plt.close()  # Free memory

    # Create a figure for the montage
    fig = plt.figure(figsize=(15, 15))  # Adjust figure size
    im_masked_noise = None
    for i in range(num_slices):
        ax = fig.add_subplot(rows, cols, i + 1)  # Create subplot for each slice
        im = ax.imshow(masked_noise[:, :, i], cmap='gray', clim=(0, 2000))  # Display slice
        if i == 0:
            im_masked_noise = im
        ax.set_title(f"Slice {i}")  # Label each slice
        ax.axis('off')  # Hide axes
    # Add colorbar on the right side
    fig.subplots_adjust(right=0.90)
    cbar_ax = fig.add_axes([0.92, 0.15, 0.02, 0.7])
    fig.colorbar(im_masked_noise, cax=cbar_ax, label='Intensity')
    # Adjust layout
    fig.tight_layout(pad=0.5, rect=[0, 0, 0.90, 1])
    ## plt.show()  # Disabled for script mode
    # Save the montage
    output_filename = 'masked_noise_volume_montage.png'
    output_path = f"{output_dir}/{output_filename}"
    plt.savefig(output_path, dpi=300)
    #plt.close()  # Free memory

    # plot iSNR SAG
    fig, axs = plt.subplots(1, 1, figsize=(12,6))  # Adjust figsize as needed
    # Plot imgm_cla
    im_cla = axs.imshow(isnr_obj_cla.isnr_map[slice_index_sag,:,:, time_point].T, origin='lower', cmap='inferno', clim=(0, isnrScale))
    axs.set_title(f'iSNR Map (Slice {slice_index})')  # Set subplot title with slice index
    axs.axis(False)  # Turn off axis labels and ticks
    cb = fig.colorbar(im_cla, ax=axs, shrink=0.3)
    cb.set_label('iSNR')
    # Adjust layout and display the plot
    plt.tight_layout()
    # plt.show()  # Disabled for script mode
    output_filename = 'iSNR_sag.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    # plot iSNR COR
    fig, axs = plt.subplots(1, 1, figsize=(12,6))  # Adjust figsize as needed
    # Plot imgm_cla
    im_cla = axs.imshow(isnr_obj_cla.isnr_map[:,slice_index_cor,:, time_point].T, origin='lower', cmap='inferno', clim=(0, isnrScale))
    axs.set_title(f'iSNR Map (Slice {slice_index})')  # Set subplot title with slice index
    axs.axis(False)  # Turn off axis labels and ticks
    cb = fig.colorbar(im_cla, ax=axs, shrink=0.3)
    cb.set_label('iSNR')
    # Adjust layout and display the plot
    plt.tight_layout()
    # plt.show()  # Disabled for script mode
    output_filename = 'iSNR_cor.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    # iSNR MONTAGE
    fig = plt.figure(figsize=(10, 10))  # Adjust figsize as needed
    im_isnr = None
    for i in range(isnr_obj_cla.isnr_map.shape[2]):
        # Create a subplot for the current slice
        ax = fig.add_subplot(rows, cols, i + 1)  # i+1 because subplot indices start from 1
        # Display the current slice using imshow
        im = ax.imshow(isnr_obj_cla.isnr_map[:, :, i, time_point], cmap='inferno', clim=(0, isnrScale))  # Adjust colormap as needed
        if i == 0:
            im_isnr = im
        ax.set_title(f"Slice {i}")  # Set title with slice index
        ax.axis('off')  # Turn off axis labels and ticks
    
    # Add colorbar on the right side
    fig.subplots_adjust(right=0.90)
    cbar_ax = fig.add_axes([0.92, 0.15, 0.02, 0.7])
    fig.colorbar(im_isnr, cax=cbar_ax, label='iSNR')
    
    # Adjust layout and spacing of subplots
    fig.tight_layout(pad=0.5, rect=[0, 0, 0.90, 1])

    # Save the montage as a PNG file
    output_filename = 'isnr_montage.png'
    output_path = f"{output_dir}/{output_filename}"
    plt.savefig(output_path, dpi=300)  # Save the montage as a PNG file with 300 dpi resolution

    ############################## tSNR
    # need to remove noise scan
    imgm_cla_nn = imgm_cla[:, :, :, :-1]  # Exclude the last volume along the time dimension

    #imgm_cla_nn = imgm_cla[:, :, :, :20]  # take first 20 dynamics.
    #imgm_cla_nn = imgm_cla #

    # Save as NIFTI 
    tsnr_obj_cla = snr.Tsnr(imgm_cla_nn, imgm_affine)
    tsnr_obj_cla.to_nifti(output_dir, 'tsnr')

    print(tsnr_obj_cla.tsnr_map.shape)

    my_mean_tsnr = np.mean(tsnr_obj_cla.tsnr_map)
    print("Mean tSNR:", my_mean_tsnr)

    # Ensure mask matches data shape; if not, ignore mask to avoid indexing errors
    if mask_data is not None:
        if getattr(mask_data, 'shape', None) is None or tuple(mask_data.shape[:3]) != tuple(tsnr_obj_cla.tsnr_map.shape[:3]):
            try:
                print(f"Mask shape {None if getattr(mask_data,'shape',None) is None else mask_data.shape} does not match data {tsnr_obj_cla.tsnr_map.shape[:3]}; skipping mask for this dataset.")
            except Exception:
                print("Mask shape mismatch; skipping mask for this dataset.")
            mask_data = None

    if mask_data is not None:

        # Thresholded tSNR inside mask
        tsnr_threshold = np.percentile(tsnr_obj_cla.tsnr_map[mask_data > 0], 50)
        voxels_to_plot = (tsnr_obj_cla.tsnr_map > tsnr_threshold) & (mask_data > 0)

        print("Found mask")
        masked_tSNR = tsnr_obj_cla.tsnr_map[mask_data > 0]
        print("Mean tSNR in mask:", np.mean(masked_tSNR))
        # Plot a tSNR slice with masking
        thisslice = tsnr_obj_cla.tsnr_map.shape[2] // 2  # Middle slice
        tsnr_slice = tsnr_obj_cla.tsnr_map[:, :, thisslice]
        mask_slice = mask_data[:, :, thisslice]

        fig, ax = plt.subplots(1, 1, figsize=(8, 8))
        im = ax.imshow(np.rot90(np.where(mask_slice > 0, tsnr_slice, np.nan)), 
                       cmap='inferno', vmin=0, vmax=np.nanmax(tsnr_slice))
        ax.set_title(f"tSNR (slice {thisslice}) in mask")
        ax.axis('off')
        cb = fig.colorbar(im, ax=ax, shrink=0.6)
        cb.set_label('tSNR')
        plt.tight_layout()
        # plt.show()  # Disabled for script mode
        fig.savefig(os.path.join(output_dir, "tSNR_masked_slice.png"), dpi=300)
        plt.close()
        
        # 3D plot
        fig = plt.figure(figsize=(10, 10))
        ax = fig.add_subplot(111, projection='3d')
        
        # Overlay masked tSNR voxels
        ax.voxels(voxels_to_plot, facecolors='orange', edgecolor='k', alpha=0.6)
        
        ax.set_title("3D tSNR voxels")
        ax.set_xlabel('X')
        ax.set_ylabel('Y')
        ax.set_zlabel('Z')
        plt.tight_layout()
        plt.savefig(os.path.join(output_dir, "tSNR_masked_3D.png"), dpi=300)
        plt.close()

    # tSNR SAG
    fig, axs = plt.subplots(1, 1, figsize=(12,6))  # Adjust figsize as needed
    # Plot imgm_cla
    im_cla = axs.imshow(tsnr_obj_cla.tsnr_map[slice_index_sag,:,:].T, origin='lower', cmap='inferno', clim=(0, tsnrScale))
    #im_cla = axs.imshow(isnr_obj_cla.isnr_map[:,:,slice_index, time_point].T, origin='lower', cmap='inferno')
    axs.set_title(f'tSNR Map (Slice {slice_index})')  # Set subplot title with slice index
    axs.axis(False)  # Turn off axis labels and ticks
    cb = fig.colorbar(im_cla, ax=axs, shrink=0.3)
    cb.set_label('tSNR')
    # Adjust layout and display the plot
    plt.tight_layout()
    # plt.show()  # Disabled for script mode
    output_filename = 'tSNR_sag.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    # tSNR COR
    fig, axs = plt.subplots(1, 1, figsize=(12,6))  # Adjust figsize as needed
    # Plot imgm_cla
    im_cla = axs.imshow(tsnr_obj_cla.tsnr_map[:,slice_index_cor,:].T, origin='lower', cmap='inferno', clim=(0, tsnrScale))
    #im_cla = axs.imshow(isnr_obj_cla.isnr_map[:,:,slice_index, time_point].T, origin='lower', cmap='inferno')
    axs.set_title(f'tSNR Map (Slice {slice_index})')  # Set subplot title with slice index
    axs.axis(False)  # Turn off axis labels and ticks
    cb = fig.colorbar(im_cla, ax=axs, shrink=0.3)
    cb.set_label('tSNR')
    # Adjust layout and display the plot
    plt.tight_layout()
    # plt.show()  # Disabled for script mode
    output_filename = 'tSNR_cor.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    #### tSNR per unit time ######
    
    # Determine TR (RepetitionTime) from multiple sources
    if TR is None and nifti_path is not None:
        # Try to read TR from JSON sidecar file
        TR = get_tr_from_json(nifti_path)
    
    if TR is None:
        # Fall back to default TR if not found in JSON
        TR = 1.4  # Default for MB3 BOLD data
        print(f"Using default TR = {TR}s (no JSON file found or RepetitionTime not specified)")
    else:
        print(f"Using TR = {TR}s from JSON metadata")
    
    # Get TE from JSON sidecar
    TE = None
    if nifti_path is not None:
        TE = get_te_from_json(nifti_path)
    
    # Get voxel size from NIfTI header
    voxel_size = None
    if nifti_img is not None:
        voxel_size = get_voxel_size(nifti_img)
    
    # Set Ernst angle scaling factor based on TR
    # These are approximate values for typical T1 relaxation times
    if TR <= 0.7:
        ErnstScaling = 0.5745  # For very short TR
    elif TR <= 1.0:
        ErnstScaling = 0.7071  # For short TR  
    elif TR <= 1.5:
        ErnstScaling = 0.8155  # For medium TR
    else:
        ErnstScaling = 1.0     # For longer TR (TR >= 2s, T1 = 2132ms)
    
    print(f"Using Ernst scaling factor = {ErnstScaling} for TR = {TR}s")

    # Compute tSNR per unit time
    tsnr_unit_time_map = (tsnr_obj_cla.tsnr_map / np.sqrt(TR)) * ErnstScaling
    
    # Save the new map as a NIfTI
    tsnr_unit_time_img = nib.Nifti1Image(tsnr_unit_time_map, affine=imgm_affine)
    nib.save(tsnr_unit_time_img, f"{output_dir}/tsnr_unit_time.nii.gz")
    
    # (Optional) Compute mean and log it
    mean_tsnr_unit_time = np.mean(tsnr_unit_time_map)
    print("Mean tSNR per unit time:", mean_tsnr_unit_time)
    
    # Plot and save image
    fig_ut, axs_ut = plt.subplots(1, 1, figsize=(8, 8))
    im_ut = axs_ut.imshow(np.rot90(tsnr_unit_time_map[:, :, slice_index]), cmap='inferno', clim=(0, tsnrScale))
    axs_ut.set_title(f'tSNR per unit time (slc {slice_index})')
    axs_ut.axis(False)
    cb_ut = fig_ut.colorbar(im_ut, ax=axs_ut, shrink=0.6)
    cb_ut.set_label('tSNR / √TR')
    fig_ut.tight_layout()
    fig_ut.savefig(f"{output_dir}/tSNR_per_unit_time.png", dpi=300)
    plt.close(fig_ut)

    # Plot and save raw tSNR image
    fig_raw, axs_raw = plt.subplots(1, 1, figsize=(8, 8))
    im_raw = axs_raw.imshow(np.rot90(tsnr_obj_cla.tsnr_map[:, :, slice_index]), cmap='inferno', clim=(0, tsnrScale))
    axs_raw.set_title(f'raw tSNR (slc {slice_index})')
    axs_raw.axis(False)
    cb_raw = fig_raw.colorbar(im_raw, ax=axs_raw, shrink=0.6)
    cb_raw.set_label('tSNR')
    fig_raw.tight_layout()
    # plt.show()  # Disabled for script mode
    fig_raw.savefig(f"{output_dir}/tSNR_raw.png", dpi=300)
    plt.close(fig_raw)

    ##############################
    
    # now montage tSNR
    # Determine grid size for montage
    num_slices = tsnr_obj_cla.tsnr_map.shape[2]
    rows = int(np.ceil(np.sqrt(num_slices)))
    cols = int(np.ceil(num_slices / rows))
    
    # Create figure with larger size for bigger individual plots
    figsize_scale = 2.5  # Scale this to make plots larger
    fig = plt.figure(figsize=(cols * figsize_scale, rows * figsize_scale))
    
    # Loop through slices
    im_tsnr = None
    for i in range(num_slices):
        ax = fig.add_subplot(rows, cols, i + 1)
        im = ax.imshow(tsnr_obj_cla.tsnr_map[:, :, i], cmap='inferno', clim=(0, tsnrScale))
        if i == 0:
            im_tsnr = im
        ax.set_title(f"Slice {i}", fontsize=8)  # Smaller font size
        ax.axis('off')
    
    # Add colorbar on the right side
    fig.subplots_adjust(right=0.90)
    cbar_ax = fig.add_axes([0.92, 0.15, 0.02, 0.7])
    fig.colorbar(im_tsnr, cax=cbar_ax, label='tSNR')
    
    # Tighter layout
    fig.tight_layout(pad=0.3, rect=[0, 0, 0.90, 1])
    
    # Save output
    output_filename = 'tSNR_montage.png'
    output_path = f"{output_dir}/{output_filename}"
    plt.savefig(output_path, dpi=300, bbox_inches='tight')

    ############################## tSNR in a patch ROI

    #slice_index = 10  # Adjust this to the desired slice index
    #time_points_index = 1  # Index of the time dimension (4th dimension) in imgm_cla
    
    # Extract the 2D slice at the specified index from the 3D image data
    slice_data = imgm_cla_nn[:, :, slice_index,:]
    #slice_data_mtx = imgm_mtx_nn[:, :, slice_index,:]
    print(f"ROI analysis: using slice {slice_index} with {slice_data.shape[2]} timepoints (shape: {slice_data.shape})")

    # Calculate the coordinates of the ROI
    x_end = x_start + roi_width
    y_end = y_start + roi_height

    # Create a Rectangle patch for the ROI on subplot axs[0]
    roi_rect_0 = patches.Rectangle((x_start, y_start), roi_width, roi_height,
                                   linewidth=1, edgecolor='y', linestyle='--', fill=False)

    # Create a 1x2 grid of subplots
    fig, axs = plt.subplots(1, 2, figsize=(8, 4))  # Adjust figsize as needed

    # Plot imgm_cla
    im_cla = axs[0].imshow(np.rot90(tsnr_obj_cla.tsnr_map[:, :, slice_index]), cmap='inferno', clim=(0, tsnrScale))
    axs[0].set_title(f'tSNR 2D Slice ROI (slc {slice_index})')  # Set the plot title
    axs[0].axis(False)  # Turn off axis labels and ticks

    # Plot the same slice on the second subplot with the rectangle
    im_gra = axs[1].imshow(np.rot90(tsnr_obj_cla.tsnr_map[:, :, slice_index]), cmap='gray')
    axs[1].add_patch(roi_rect_0)
    axs[1].set_title(f'Slice {slice_index} with ROI')  # Set the plot title
    axs[1].axis(False)  # Turn off axis labels and ticks
    
    cb = fig.colorbar(im_cla, ax=axs[0],shrink=0.6) #, fraction=0.046, pad=0.04)
    cb.set_label('Intensity')
    cb2 = fig.colorbar(im_gra, ax=axs[1],shrink=0.6) #, fraction=0.046, pad=0.04)
    #cb.set_label('Intensity')
    #fig.colorbar(im, location='bottom')

    # Adjust layout and display the plot
    plt.tight_layout()
    # plt.show()  # Disabled for script mode

    output_filename = 'tSNR_w_ROI_images.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    ############################## signal and std over time

    # Crop the selected ROI from the 2D slice data
    roi_data = slice_data[y_start:y_end, x_start:x_end, :]

    # Calculate the average signal intensity across the selected ROI over time
    average_patch = np.mean(roi_data, axis=(0, 1))

    # Calculate the standard deviation (STDev) of the signal intensity across the selected ROI over time
    std_patch = np.std(roi_data, axis=(0, 1))

    # Ensure no NaNs or Infs before detrending to avoid SciPy errors
    if not np.all(np.isfinite(std_patch)):
        std_patch = np.nan_to_num(std_patch, nan=0.0, posinf=0.0, neginf=0.0)
    try:
        std_patch = detrend(std_patch)
    except Exception:
        # If detrend still fails, fall back to the raw (sanitized) series
        pass

    # Detrend the average patch time series, with the same safety handling
    if not np.all(np.isfinite(average_patch)):
        average_patch = np.nan_to_num(average_patch, nan=0.0, posinf=0.0, neginf=0.0)
    try:
        detrended_patch = detrend(average_patch)
    except Exception:
        detrended_patch = average_patch

    # Prepare the time points (x-axis) for the time series plot
    time_points = np.arange(slice_data.shape[2])

    #Plot the time series of the average patch
    fig = plt.figure(figsize=(8, 4))  # Adjust figsize as needed
    plt.plot(time_points, detrended_patch, color='blue', label='Average Patch')

    #Plot the detrended time series of the average patch_mtx (second time series)
    plt.plot(time_points, std_patch, color='red', label='stdev Patch')

    #Set plot title and labels
    plt.title(f"Time Series of patch (Slice {slice_index})")
    plt.xlabel("Time (Index)")
    plt.ylabel("Signal Intensity")
    plt.legend()  # Show legend with labels for each time series
    #plt.ylim(-400,400)
    #Show the plot
    plt.grid(True)  # Enable grid for better visualization
    # plt.show()  # Disabled for script mode
    output_filename = 'TS_images.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory

    tsnr_slice = tsnr_obj_cla.tsnr_map[:, :, slice_index]
    tsnr_slice = np.rot90(tsnr_slice)  # Rotate the slice if needed
    #print("Shape of tSNR slice:", tsnr_slice.shape)
    tsnr_roi = tsnr_slice[y_start:y_start + roi_height, x_start:x_start + roi_width]
    mean_tsnr_roi = np.mean(tsnr_roi)
    print("Mean tSNR within ROI:", mean_tsnr_roi)

    
    ############################## STATIC spatial noise image

    # Lastly, let's plot the static spatial noise images.

    slice_data_odd = imgm_cla_nn[:, :, :,::2]
    slice_data_even = imgm_cla_nn[:, :, :,1::2]

    # Sum across the fourth dimension (time) to get the sum of odd and even slices
    sum_odd = np.sum(slice_data_odd, axis=3)
    sum_even = np.sum(slice_data_even, axis=3)

    # Calculate the difference between sum of odd and even slices
    static_spatial_noise = sum_odd - sum_even

    # Apply thresholding to remove background noise
    threshold_value = 0  # Adjust threshold value as needed
    static_spatial_noise_thresholded = np.where(static_spatial_noise < threshold_value, 0, static_spatial_noise)

    # Select a specific slice index (e.g., quickCrop(5) in MATLAB)
    slice_index = 5  # Adjust as needed

    # Plot the static spatial noise image
    fig = plt.figure(figsize=(6, 4))  # Adjust figsize as needed
    plt.imshow(static_spatial_noise[:, :, slice_index], cmap='viridis', aspect='equal')
    plt.title(f"Static Spatial Noise Image, mean={int(np.round(np.mean(static_spatial_noise)))}")
    plt.colorbar(label='Intensity')

    # Set colorbar limits (clim) if desired
    #plt.clim(-100, 100)

    # Show the plot
    # plt.show()  # Disabled for script mode
    output_filename = 'SSN.png'
    output_path = f"{output_dir}/{output_filename}"  # Construct the full output path
    fig.savefig(output_path, dpi=300)  # Save the plot as a PNG file with 300 dpi resolution
    plt.close()  # Close the plot to free up memory
    mean_ssn = np.mean(static_spatial_noise)
    print("Mean SSN:", mean_ssn)

    # Collect all metrics in a dictionary
    metrics = {
        'filename': core_filename,
        'num_timepoints': imgm_cla_nn.shape[3],
        'num_slices': imgm_cla.shape[2],
        'slice_index_analyzed': slice_index,
        'mean_volume_std': mean_std,
        'mean_isnr': np.mean(isnr_obj_cla.isnr),
        'noise_value': np.mean(isnr_obj_cla.noise),
        'mean_tsnr': my_mean_tsnr,
        'mean_tsnr_in_mask': np.mean(masked_tSNR) if mask_data is not None else None,
        'TR': TR,
        'TE': TE,
        'voxel_size': voxel_size,
        'ernst_scaling': ErnstScaling,
        'mean_tsnr_per_unit_time': mean_tsnr_unit_time,
        'mean_tsnr_roi': mean_tsnr_roi,
        'mean_ssn': mean_ssn,
        'roi_slice': slice_index,
        'roi_x_start': x_start,
        'roi_y_start': y_start,
        'roi_width': roi_width,
        'roi_height': roi_height
    }

    # The End
    return metrics

def create_qa_powerpoint(output_dir, subject_name="QA Analysis"):
    """
    Create a PowerPoint presentation from QA output images
    
    Parameters:
    -----------
    output_dir : str
        Path to the QA output directory containing PNG images
    subject_name : str
        Name to use in the title slide
    """
    
    if not PPTX_AVAILABLE:
        print("Skipping PowerPoint generation (python-pptx not installed)")
        return None
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define image order and titles
    # Order: Mean/Noise, then tSNR (all datasets), then iSNR (all datasets), then other analyses
    image_config = [
        # Mean and noise images first
        ('Mean_image.png', 'Mean Image'),
        ('mean_montage.png', 'Mean Image Montage - All Slices'),
        ('masked_noise.png', 'Noise Volume Analysis'),
        ('noise_volume_montage.png', 'Noise Volume Montage'),
        ('masked_noise_volume_montage.png', 'Masked Noise Volume Montage'),
        # tSNR images (all datasets show tSNR first)
        ('tSNR_sag.png', 'tSNR Map - Sagittal View'),
        ('tSNR_cor.png', 'tSNR Map - Coronal View'),
        ('tSNR_per_unit_time.png', 'tSNR per Unit Time'),
        ('tSNR_raw.png', 'Raw tSNR Map'),
        ('tSNR_montage.png', 'tSNR Montage - All Slices'),
        ('tSNR_w_ROI_images.png', 'tSNR with ROI'),
        # iSNR images (all datasets show iSNR after tSNR)
        ('iSNR_sag.png', 'iSNR Map - Sagittal View'),
        ('iSNR_cor.png', 'iSNR Map - Coronal View'),
        ('isnr_montage.png', 'iSNR Montage - All Slices'),
        # Other analyses
        ('TS_images.png', 'Time Series Analysis'),
        ('SSN.png', 'Static Spatial Noise'),
    ]
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assurance Report"
    subtitle.text = f"{subject_name}\n{os.path.basename(output_dir)}"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Add slides for each image
    for img_filename, img_title in image_config:
        img_path = os.path.join(output_dir, img_filename)
        
        if not os.path.exists(img_path):
            continue
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = img_title
        
        # Style title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Calculate image dimensions to fit on slide
        # Available space: slide height - title height - margins
        max_width = Inches(9)
        max_height = Inches(6.2)  # 7.5 - 0.3 (top) - 0.6 (title) - 0.4 (margin)
        
        # Get image dimensions
        try:
            img = Image.open(img_path)
            img_width, img_height = img.size
            
            # Calculate scaling to fit within max dimensions while preserving aspect ratio
            width_ratio = max_width / Inches(img_width / 96)  # Assuming 96 DPI
            height_ratio = max_height / Inches(img_height / 96)
            scale_ratio = min(width_ratio, height_ratio, 1.0)  # Don't upscale
            
            final_width = Inches(img_width / 96) * scale_ratio
            final_height = Inches(img_height / 96) * scale_ratio
            
            # Center the image horizontally
            img_left = (Inches(10) - final_width) / 2
            img_top = Inches(1.1)
            
            slide.shapes.add_picture(img_path, img_left, img_top, width=final_width, height=final_height)
        except Exception as e:
            # Fallback: use fixed width
            img_left = Inches(0.5)
            img_top = Inches(1.1)
            slide.shapes.add_picture(img_path, img_left, img_top, width=max_width)
    
    # Add CSV metrics slide
    csv_files = glob(os.path.join(os.path.dirname(output_dir), 'qa_metrics_*.csv'))
    if csv_files:
        csv_path = csv_files[-1]  # Get most recent
        
        # Create blank slide for table
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        tf.text = "QA Metrics Summary"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Read CSV and add as table
        try:
            with open(csv_path, 'r') as f:
                reader = csv.DictReader(f)
                rows = list(reader)
                
            if rows:
                # Create table with key metrics only (to fit on slide)
                key_metrics = [
                    'num_timepoints', 'num_slices', 'mean_volume_std',
                    'mean_isnr', 'mean_tsnr', 'TR', 'mean_tsnr_per_unit_time',
                    'mean_tsnr_roi', 'mean_ssn'
                ]
                
                # Filter to key metrics that exist
                available_metrics = [m for m in key_metrics if m in rows[0]]
                
                # Create table (2 columns: metric name, value)
                table_rows = len(available_metrics) + 1  # +1 for header
                table_cols = 2
                
                left = Inches(1.5)
                top = Inches(1.2)
                width = Inches(7)
                height = Inches(5.5)
                
                table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
                
                # Set column widths
                table.columns[0].width = Inches(4)
                table.columns[1].width = Inches(3)
                
                # Header row
                table.cell(0, 0).text = "Metric"
                table.cell(0, 1).text = "Value"
                for cell in [table.cell(0, 0), table.cell(0, 1)]:
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(12)
                
                # Data rows
                for i, metric in enumerate(available_metrics):
                    # Format metric name
                    metric_display = metric.replace('_', ' ').title()
                    table.cell(i + 1, 0).text = metric_display
                    
                    # Format value
                    value = rows[0][metric]
                    try:
                        # Try to format as float with 2 decimals
                        float_val = float(value)
                        table.cell(i + 1, 1).text = f"{float_val:.2f}"
                    except (ValueError, TypeError):
                        table.cell(i + 1, 1).text = str(value)
                    
                    # Set font size
                    for j in range(2):
                        table.cell(i + 1, j).text_frame.paragraphs[0].font.size = Pt(11)
                        
        except Exception as e:
            print(f"Could not add CSV table: {e}")
    
    # Save presentation
    pptx_filename = os.path.join(output_dir, f"QA_Report_{os.path.basename(output_dir)}.pptx")
    prs.save(pptx_filename)
    print(f"\nPowerPoint created: {pptx_filename}")
    
    return pptx_filename

def run_qa_multi_datasets(base_data_dir, extension='.nii.gz', filename_pattern='*task-rest-bold', 
                          base_output_dir='/Users/cmilbourn/Documents/tSNR_check'):
    """
    Run QA analysis for multiple datasets and combine into single PowerPoint and CSV
    
    Parameters:
    -----------
    base_data_dir : str
        Base directory containing multiple subject/session folders
    extension : str
        File extension (.nii.gz or .nii)
    filename_pattern : str
        Pattern to match fMRI files
    base_output_dir : str
        Base directory for all QA outputs
    """
    
    # Find all fMRI files matching the pattern
    search_pattern = os.path.join(base_data_dir, '**', f'{filename_pattern}{extension}')
    all_files = glob(search_pattern, recursive=True)
    
    print(f"Found {len(all_files)} files matching pattern:")
    for f in all_files:
        print(f"  {f}")
    
    if not all_files:
        print("No files found!")
        return
    
    # Create base output directory
    os.makedirs(base_output_dir, exist_ok=True)
    
    # Create a timestamped session directory
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    session_dir = os.path.join(base_output_dir, f'qa_session_{timestamp}')
    os.makedirs(session_dir, exist_ok=True)
    
    # List to store all metrics and output directories
    all_metrics = []
    all_output_dirs = []
    # Tracking arrays for per-file mask handling (initialized in main scope)
    created_masks = []            # masks found pre-existing with expected naming
    generated_masks = []          # masks created via BET during this run
    missing_mask_failures = []    # datasets where mask creation or loading failed
    # Tracking arrays for mask handling (main scope)
    created_masks = []
    generated_masks = []
    missing_mask_failures = []
    # Tracking for mask handling (defined in main scope)
    created_masks = []          # existing masks matching filename
    generated_masks = []        # masks generated via BET
    missing_mask_failures = []  # datasets where BET failed or mask load failed
    # Track mask usage/creation
    created_masks = []          # masks found and used
    generated_masks = []        # masks generated via BET
    missing_mask_failures = []  # datasets where BET failed
    
    # Process each file
    for file_idx, mag_file_path in enumerate(all_files):
        print(f"\n{'='*80}")
        print(f"Processing file {file_idx + 1}/{len(all_files)}: {os.path.basename(mag_file_path)}")
        print(f"{'='*80}")
        
        # Extract a clean name for this dataset
        core_filename = os.path.splitext(os.path.splitext(os.path.basename(mag_file_path))[0])[0]
        pathname_m = os.path.dirname(mag_file_path)

        # Extract a clean name for this dataset
        core_filename = os.path.splitext(os.path.splitext(os.path.basename(mag_file_path))[0])[0]
        pathname_m = os.path.dirname(mag_file_path)

        print(f"Core filename: {core_filename}")

        # Create an output directory for saving plots
        output_directory = os.path.join(session_dir, f'{core_filename}')
        os.makedirs(output_directory, exist_ok=True)
        OUTPUT_DIR = os.path.abspath(output_directory)

        print(f"Output directory: {OUTPUT_DIR}")
        
        # Remove first 2 dummy volumes using fslroi
        mag_file_nodummies = os.path.join(OUTPUT_DIR, core_filename + '_nodummies' + extension)
        print(f"Removing first 2 dummy volumes with fslroi...")
        fslroi_cmd = f"fslroi {mag_file_path} {mag_file_nodummies} 2 -1"
        result = subprocess.run(fslroi_cmd, shell=True, capture_output=True, text=True)
        if result.returncode != 0:
            print(f"Warning: fslroi failed: {result.stderr}")
            print("Proceeding with original file...")
            mag_file_to_use = mag_file_path
        else:
            print(f"Dummy volumes removed. Using: {mag_file_nodummies}")
            mag_file_to_use = mag_file_nodummies
        
        # Load magnitude data only
        print('Loading magnitude data...')
        imgm_cla, imgm_cla_affine = load_data(mag_file_to_use)

        # MASK (prefer per-file mask matching this magnitude)
        mag_base = os.path.splitext(os.path.splitext(os.path.basename(mag_file_to_use))[0])[0]
        mask_path = find_mask_file(pathname_m, mag_basename=mag_base)
        if mask_path:
            print(f"Found mask: {mask_path}")
            mask_data, mask_affine = load_data(mask_path)
        else:
            print("No mask file found.")
            mask_data = None
        
        # Process and plot data
        metrics = process_data_nophase(imgm_cla, imgm_cla_affine, core_filename, OUTPUT_DIR, mask_data, nifti_path=mag_file_path)
        all_metrics.append(metrics)
        all_output_dirs.append(OUTPUT_DIR)
    
    # Create combined PowerPoint presentation
    print(f"\n{'='*80}")
    print("Creating combined PowerPoint presentation...")
    print(f"{'='*80}")
    create_combined_powerpoint(session_dir, all_output_dirs, all_metrics)
    
    # Save combined metrics to CSV
    if all_metrics:
        csv_path = os.path.join(session_dir, f'qa_metrics_combined_{timestamp}.csv')
        with open(csv_path, 'w', newline='') as csvfile:
            fieldnames = all_metrics[0].keys()
            writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
            writer.writeheader()
            for metrics in all_metrics:
                writer.writerow(metrics)
        print(f"\nCombined metrics CSV saved to: {csv_path}")

def create_combined_powerpoint(session_dir, all_output_dirs, all_metrics):
    """
    Create a single PowerPoint with all datasets
    
    Parameters:
    -----------
    session_dir : str
        Session directory for output
    all_output_dirs : list
        List of output directories for each dataset
    all_metrics : list
        List of metrics dictionaries for each dataset
    """
    
    if not PPTX_AVAILABLE:
        print("Skipping PowerPoint generation (python-pptx not installed)")
        return None
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define image order and titles
    image_config = [
        # Mean and noise images first
        ('Mean_image.png', 'Mean Image'),
        ('mean_montage.png', 'Mean Image Montage - All Slices'),
        ('masked_noise.png', 'Noise Volume Analysis'),
        ('noise_volume_montage.png', 'Noise Volume Montage'),
        ('masked_noise_volume_montage.png', 'Masked Noise Volume Montage'),
        # tSNR images (all datasets show tSNR first)
        ('tSNR_sag.png', 'tSNR Map - Sagittal View'),
        ('tSNR_cor.png', 'tSNR Map - Coronal View'),
        ('tSNR_per_unit_time.png', 'tSNR per Unit Time'),
        ('tSNR_raw.png', 'Raw tSNR Map'),
        ('tSNR_montage.png', 'tSNR Montage - All Slices'),
        ('tSNR_w_ROI_images.png', 'tSNR with ROI'),
        # iSNR images (all datasets show iSNR after tSNR)
        ('iSNR_sag.png', 'iSNR Map - Sagittal View'),
        ('iSNR_cor.png', 'iSNR Map - Coronal View'),
        ('isnr_montage.png', 'iSNR Montage - All Slices'),
        # Other analyses
        ('TS_images.png', 'Time Series Analysis'),
        ('SSN.png', 'Static Spatial Noise'),
    ]
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assurance Report"
    subtitle.text = f"Combined Report - {len(all_output_dirs)} Datasets\n{datetime.now().strftime('%Y-%m-%d %H:%M')}"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Add slides for each dataset
    for dataset_idx, output_dir in enumerate(all_output_dirs):
        dataset_name = os.path.basename(output_dir)
        
        # Add dataset separator slide
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.text = f"Dataset {dataset_idx + 1}/{len(all_output_dirs)}\n{dataset_name}"
        p = tf.paragraphs[0]
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add image slides for this dataset
        for img_filename, img_title in image_config:
            img_path = os.path.join(output_dir, img_filename)
            
            if not os.path.exists(img_path):
                continue
            
            # Create blank slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title with dataset name
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            tf = txBox.text_frame
            tf.text = f"{img_title} - {dataset_name}"
            
            # Style title
            p = tf.paragraphs[0]
            p.font.size = Pt(16)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            
            # For tSNR_montage only, add a table with acquisition parameters (TR, TE, voxel size, tSNR in mask)
            if img_filename == 'tSNR_montage.png' and dataset_idx < len(all_metrics):
                # Get metrics for this dataset
                dataset_metrics = all_metrics[dataset_idx]
                
                # Create table on the right side: 2 columns, 5 rows (header + 4 data rows)
                table_left = Inches(7.0)
                table_top = Inches(1.0)
                table_width = Inches(2.5)
                table_height = Inches(2.0)
                
                table = slide.shapes.add_table(5, 2, table_left, table_top, table_width, table_height).table
                
                # Set column widths
                table.columns[0].width = Inches(1.2)
                table.columns[1].width = Inches(1.3)
                
                # Header row
                table.cell(0, 0).text = 'Parameter'
                table.cell(0, 1).text = 'Value'
                
                # TR row
                table.cell(1, 0).text = 'TR'
                tr_val = dataset_metrics.get('TR')
                table.cell(1, 1).text = f"{tr_val:.3f} s" if tr_val is not None else 'N/A'
                
                # TE row
                table.cell(2, 0).text = 'TE'
                te_val = dataset_metrics.get('TE')
                if te_val is not None:
                    table.cell(2, 1).text = f"{te_val*1000:.2f} ms"
                else:
                    table.cell(2, 1).text = 'N/A'
                
                # Voxel size row
                table.cell(3, 0).text = 'Voxel Size'
                voxel_val = dataset_metrics.get('voxel_size')
                if voxel_val is not None:
                    table.cell(3, 1).text = f"{voxel_val[0]:.1f}×{voxel_val[1]:.1f}×{voxel_val[2]:.1f} mm"
                else:
                    table.cell(3, 1).text = 'N/A'
                
                # tSNR in mask row
                table.cell(4, 0).text = 'tSNR in mask'
                tsnr_mask_val = dataset_metrics.get('mean_tsnr_in_mask')
                if tsnr_mask_val is not None:
                    table.cell(4, 1).text = f"{tsnr_mask_val:.2f}"
                else:
                    table.cell(4, 1).text = 'N/A'
                
                # Style table cells
                for row_idx in range(5):
                    for col_idx in range(2):
                        cell = table.cell(row_idx, col_idx)
                        cell.text_frame.paragraphs[0].font.size = Pt(9)
                        if row_idx == 0:  # Header row
                            cell.text_frame.paragraphs[0].font.bold = True
                        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                        cell.vertical_anchor = 1  # Middle alignment
            
            # Calculate image dimensions to fit on slide
            # Adjust width for montage (with table); all others full width
            if img_filename == 'tSNR_montage.png':
                max_width = Inches(6.5)  # Leave room for table on right
            else:
                max_width = Inches(9)
            max_height = Inches(6.2)
            
            # Get image dimensions and scale
            try:
                img = Image.open(img_path)
                img_width, img_height = img.size
                
                width_ratio = max_width / Inches(img_width / 96)
                height_ratio = max_height / Inches(img_height / 96)
                scale_ratio = min(width_ratio, height_ratio, 1.0)
                
                final_width = Inches(img_width / 96) * scale_ratio
                final_height = Inches(img_height / 96) * scale_ratio
                
                # Left-align montage (table present), center others
                if img_filename == 'tSNR_montage.png':
                    img_left = Inches(0.5)
                else:
                    img_left = (Inches(10) - final_width) / 2
                img_top = Inches(1.1)
                
                slide.shapes.add_picture(img_path, img_left, img_top, width=final_width, height=final_height)
            except Exception as e:
                img_left = Inches(0.5)
                img_top = Inches(1.1)
                slide.shapes.add_picture(img_path, img_left, img_top, width=max_width)
    
    # Add combined CSV metrics table slide
    if all_metrics:
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        tf = txBox.text_frame
        tf.text = "QA Metrics Summary - All Datasets"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Create table with key metrics
        try:
            key_metrics = [
                'filename', 'num_timepoints', 'mean_isnr', 'mean_tsnr', 
                'mean_tsnr_per_unit_time', 'mean_tsnr_roi'
            ]
            
            # Filter to available metrics
            available_metrics = [m for m in key_metrics if m in all_metrics[0]]
            
            # Create table
            table_rows = len(all_metrics) + 1  # +1 for header
            table_cols = len(available_metrics)
            
            # Adjust table size based on number of datasets
            left = Inches(0.3)
            top = Inches(1.1)
            width = Inches(9.4)
            height = Inches(6)
            
            table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table
            
            # Header row
            for col_idx, metric in enumerate(available_metrics):
                metric_display = metric.replace('_', ' ').title()
                if metric == 'filename':
                    metric_display = 'Dataset'
                table.cell(0, col_idx).text = metric_display
                table.cell(0, col_idx).text_frame.paragraphs[0].font.bold = True
                table.cell(0, col_idx).text_frame.paragraphs[0].font.size = Pt(10)
            
            # Data rows
            for row_idx, metrics in enumerate(all_metrics):
                for col_idx, metric in enumerate(available_metrics):
                    value = metrics.get(metric, '')
                    if metric == 'filename':
                        # Shorten filename for display
                        value = str(value)[:40] + '...' if len(str(value)) > 40 else str(value)
                    else:
                        try:
                            float_val = float(value)
                            value = f"{float_val:.2f}"
                        except (ValueError, TypeError):
                            value = str(value)
                    
                    table.cell(row_idx + 1, col_idx).text = value
                    table.cell(row_idx + 1, col_idx).text_frame.paragraphs[0].font.size = Pt(9)
                    
        except Exception as e:
            print(f"Could not add CSV table: {e}")
    
    # Save presentation
    pptx_filename = os.path.join(session_dir, f"QA_Report_Combined_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    prs.save(pptx_filename)
    print(f"\nCombined PowerPoint created: {pptx_filename}")
    
    return pptx_filename

if __name__ == "__main__":
    # Explicit dataset list (run even if 3D/short)
    dataset_configs = [
        {
            'path': '/Users/cmilbourn/Documents/tSNR_check_30dyn/Data2run/rawish/sub005/sub005-visit001_17976-012_fmri_MB1_ARC2_fMRI_2mm_Sweet_20251125_bold_dev_20251216143719_5.nii.gz',
            'TR': None
        }
    ]

    # ------------------------------------------------------------------
    # Previous auto-discovery block (kept for reference, now disabled)
    # ------------------------------------------------------------------
    # target_dir = '/Users/cmilbourn/Documents/tSNR_check_30dyn/Data2run/rawish/sub005'
    # candidate_files = sorted(glob(os.path.join(target_dir, '*.nii*')))
    # dataset_configs = []
    #
    # for fpath in candidate_files:
    #     if not fpath.endswith(('.nii', '.nii.gz')):
    #         continue
    #     try:
    #         img = nib.load(fpath)
    #         # Only keep 4D datasets with at least 10 volumes (exclude anatomical/localizer)
    #         if img.ndim < 4 or img.shape[-1] < 10:
    #             print(f"Skipping (not 4D or too few volumes): {os.path.basename(fpath)} shape={img.shape}")
    #             continue
    #     except Exception as e:
    #         print(f"Skipping (load error) {fpath}: {e}")
    #         continue
    #     dataset_configs.append({ 'path': fpath, 'TR': None })


    # Check which files exist
    existing_configs = []
    for cfg in dataset_configs:
        fpath = cfg['path']
        if os.path.exists(fpath):
            existing_configs.append(cfg)
            print(f"Found: {os.path.basename(fpath)}")
        else:
            print(f"Path does not exist: {fpath}")

    if not existing_configs:
        print("No files found to process!")
        exit(1)

    # Replace with verified list
    dataset_configs = existing_configs
    
    print(f"\nTotal files to process: {len(dataset_configs)}")
    
    # Create base output directory
    base_output_dir = '/Users/cmilbourn/Documents/tSNR_check_30dyn/qa_output'
    os.makedirs(base_output_dir, exist_ok=True)
    
    # Create a timestamped session directory
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    session_dir = os.path.join(base_output_dir, f'qa_session_{timestamp}')
    os.makedirs(session_dir, exist_ok=True)
    
    # List to store all metrics and output directories
    all_metrics = []
    all_output_dirs = []
    
    # Tracking arrays for per-file mask handling
    created_masks = []            # masks found pre-existing with expected naming
    generated_masks = []          # masks created via BET during this run
    missing_mask_failures = []    # datasets where mask creation or loading failed
    
    # Per-file processing log for end-of-run summary
    processing_log = []
    
    # Process each file
    for file_idx, cfg in enumerate(dataset_configs):
        mag_file_path = cfg['path']
        tr_override = cfg.get('TR')
        print(f"\n{'='*80}")
        print(f"Processing file {file_idx + 1}/{len(dataset_configs)}: {os.path.basename(mag_file_path)}")
        print(f"{'='*80}")
        
        try:
            # Extract a clean name for this dataset
            core_filename = os.path.splitext(os.path.splitext(os.path.basename(mag_file_path))[0])[0]
            pathname_m = os.path.dirname(mag_file_path)

            print(f"Core filename: {core_filename}")

            # Create an output directory for saving plots
            output_directory = os.path.join(session_dir, f'{core_filename}')
            os.makedirs(output_directory, exist_ok=True)
            OUTPUT_DIR = os.path.abspath(output_directory)

            print(f"Output directory: {OUTPUT_DIR}")
            if tr_override is not None:
                print(f"Using TR override = {tr_override}s for this dataset")
            
            # Remove first 2 dummy volumes using fslroi (only if 4D with >=3 volumes)
            mag_file_to_use = mag_file_path
            try:
                _hdr_probe = nib.load(mag_file_path)
                is_4d = (_hdr_probe.ndim == 4)
                t_len = _hdr_probe.shape[-1] if is_4d else 1
            except Exception as e:
                print(f"Warning: could not probe NIfTI header, skipping dummy removal: {e}")
                is_4d = False
                t_len = 1

            if is_4d and t_len >= 3:
                mag_file_nodummies = os.path.join(OUTPUT_DIR, core_filename + '_nodummies.nii.gz')
                print(f"Removing first 2 dummy volumes with fslroi...")
                fslroi_cmd = f"fslroi {mag_file_path} {mag_file_nodummies} 2 -1"
                result = subprocess.run(fslroi_cmd, shell=True, capture_output=True, text=True)
                if result.returncode != 0 or (not os.path.exists(mag_file_nodummies)):
                    print(f"Warning: fslroi failed or output missing: {result.stderr}")
                    print("Proceeding with original file...")
                    mag_file_to_use = mag_file_path
                else:
                    print(f"Dummy volumes removed. Using: {mag_file_nodummies}")
                    mag_file_to_use = mag_file_nodummies
            else:
                print(f"Skipping dummy removal (not 4D or too few timepoints: ndim={4 if is_4d else 3}, t={t_len})")

            # Crop/rename to 30 dynamics if 4D
            try:
                _hdr_probe2 = nib.load(mag_file_to_use)
                is_4d_post = (_hdr_probe2.ndim == 4)
                t_len_post = _hdr_probe2.shape[-1] if is_4d_post else 1
            except Exception as e:
                print(f"Warning: could not probe NIfTI header after dummy removal: {e}")
                is_4d_post = False
                t_len_post = 1

            if is_4d_post and t_len_post >= 1:
                keep_vols = min(30, t_len_post)
                mag_file_30dyn = os.path.join(OUTPUT_DIR, core_filename + '_30dyn.nii.gz')
                print(f"Cropping to first {keep_vols} dynamics with fslroi...")
                fslroi_cmd_30 = f"fslroi {mag_file_to_use} {mag_file_30dyn} 0 {keep_vols}"
                result30 = subprocess.run(fslroi_cmd_30, shell=True, capture_output=True, text=True)
                if result30.returncode != 0 or (not os.path.exists(mag_file_30dyn)):
                    print(f"Warning: fslroi 30dyn failed or output missing: {result30.stderr}")
                    print("Proceeding with current file (no 30dyn crop)...")
                else:
                    print(f"30-dynamic file created. Using: {mag_file_30dyn}")
                    mag_file_to_use = mag_file_30dyn
            else:
                print(f"Skipping 30dyn crop (not 4D or no time dimension: ndim={4 if is_4d_post else 3}, t={t_len_post})")
            
            # Load magnitude data only
            print('Loading magnitude data...')
            imgm_cla, imgm_cla_affine = load_data(mag_file_to_use)
            
            # Also load the NIfTI image object to get header info
            nifti_img = nib.load(mag_file_to_use)

            # MASK HANDLING (per-dataset) -------------------------------------------------
            # Search for existing _brain file matching this magnitude file
            mask_data = None
            mag_base = os.path.splitext(os.path.splitext(os.path.basename(mag_file_to_use))[0])[0]
            
            # Expect mask to match input file name with _brain_mask appended before extension
            if mag_file_path.endswith('.nii.gz'):
                expected_mask = mag_file_path.replace('.nii.gz', '_brain_mask.nii.gz')
            elif mag_file_path.endswith('.nii'):
                expected_mask = mag_file_path.replace('.nii', '_brain_mask.nii.gz')
            else:
                expected_mask = mag_file_path + '_brain_mask.nii.gz'
            
            if os.path.exists(expected_mask):
                print(f"Using existing per-file mask: {expected_mask}")
                try:
                    mask_data, mask_affine = load_data(expected_mask)
                    created_masks.append(expected_mask)
                    print(f"Successfully loaded mask from: {expected_mask}")
                except Exception as me:
                    print(f"Warning: could not load existing mask {expected_mask}: {me}")
                    mask_data = None
            else:
                # No matching mask found - try BET to generate one
                print(f"No matching mask found. Creating mask with BET for {core_filename} ...")
                bet_output_stem = os.path.join(pathname_m, mag_base + '_brain')  # Create in same dir as data
                bet_cmd = f"bet {mag_file_to_use} {bet_output_stem} -f 0.5 -m"
                bet_res = subprocess.run(bet_cmd, shell=True, capture_output=True, text=True)
                if bet_res.returncode != 0:
                    print(f"BET failed for {core_filename}: {bet_res.stderr}")
                    mask_data = None
                    missing_mask_failures.append(core_filename)
                else:
                    # BET generates bet_output_stem.nii.gz and bet_output_stem_mask.nii.gz
                    generated_mask_source = bet_output_stem + '_mask.nii.gz'
                    if os.path.exists(generated_mask_source):
                        try:
                            # Copy to expected location for future runs
                            if generated_mask_source != expected_mask:
                                try:
                                    import shutil
                                    shutil.copy2(generated_mask_source, expected_mask)
                                    print(f"Copied generated mask to: {expected_mask}")
                                except Exception as ce:
                                    print(f"Warning: could not copy mask to expected name: {ce}")
                            mask_data, mask_affine = load_data(generated_mask_source)
                            generated_masks.append(generated_mask_source)
                            print(f"Mask generated: {generated_mask_source}")
                        except Exception as me:
                            print(f"Warning: generated mask could not be loaded {generated_mask_source}: {me}")
                            mask_data = None
                            missing_mask_failures.append(core_filename)
                    else:
                        print(f"BET reported success but mask file missing: {generated_mask_source}")
                        mask_data = None
                        missing_mask_failures.append(core_filename)
            
            # Process and plot data
            metrics = process_data_nophase(imgm_cla, imgm_cla_affine, core_filename, OUTPUT_DIR, mask_data, TR=tr_override, nifti_path=mag_file_path, nifti_img=nifti_img)
            all_metrics.append(metrics)
            all_output_dirs.append(OUTPUT_DIR)
            # Record successful processing in log
            mask_path_used = expected_mask if mask_data is not None and os.path.exists(expected_mask) else None
            processing_log.append({
                'input': mag_file_path,
                'status': 'completed',
                'output_dir': OUTPUT_DIR,
                'mask_used': bool(mask_data is not None),
                'mask_path': mask_path_used,
                'TR': metrics.get('TR')
            })
            
        except Exception as e:
            print(f"ERROR processing {mag_file_path}: {e}")
            import traceback
            traceback.print_exc()
            # Record failure in log
            processing_log.append({
                'input': mag_file_path,
                'status': 'failed',
                'error': str(e)
            })
            continue
    
    # Create combined PowerPoint presentation
    print(f"\n{'='*80}")
    print("Creating combined PowerPoint presentation...")
    print(f"{'='*80}")
    create_combined_powerpoint(session_dir, all_output_dirs, all_metrics)
    
    # Save combined metrics to CSV
    if all_metrics:
        csv_path = os.path.join(session_dir, f'qa_metrics_combined_{timestamp}.csv')
        with open(csv_path, 'w', newline='') as csvfile:
            fieldnames = all_metrics[0].keys()
            writer = csv.DictWriter(csvfile, fieldnames=fieldnames)
            writer.writeheader()
            for metrics in all_metrics:
                writer.writerow(metrics)
        print(f"\nCombined metrics CSV saved to: {csv_path}")

    # Summary of mask handling
    print("\nMask processing summary:")
    print(f"  Existing masks used: {len(created_masks)}")
    print(f"  Masks generated via BET: {len(generated_masks)}")
    if missing_mask_failures:
        print(f"  Datasets with mask creation failures: {len(missing_mask_failures)}")
        for fail in missing_mask_failures:
            print(f"    - {fail}")
    else:
        print("  No mask creation failures.")
    
    # Per-file processing summary list
    print("\nProcessing results by input file:")
    completed_count = sum(1 for r in processing_log if r.get('status') == 'completed')
    failed_count = sum(1 for r in processing_log if r.get('status') == 'failed')
    total_requested = len(dataset_configs)
    print(f"  Requested: {total_requested}")
    print(f"  Completed: {completed_count}")
    print(f"  Failed: {failed_count}")
    
    for rec in processing_log:
        if rec.get('status') == 'completed':
            mused = 'yes' if rec.get('mask_used') else 'no'
            tr_val = rec.get('TR')
            print(f"    [OK] {os.path.basename(rec['input'])} | TR={tr_val} | mask_used={mused}")
        else:
            print(f"    [FAIL] {os.path.basename(rec['input'])} | error={rec.get('error')}")
    
    # Write processing log CSV
    log_csv = os.path.join(session_dir, 'processing_log.csv')
    with open(log_csv, 'w', newline='') as f:
        fieldnames = ['input', 'status', 'output_dir', 'mask_used', 'mask_path', 'TR', 'error']
        writer = csv.DictWriter(f, fieldnames=fieldnames)
        writer.writeheader()
        for rec in processing_log:
            writer.writerow(rec)
    print(f"\nProcessing log saved to: {log_csv}")
    
    print(f"\n{'='*80}")
    print(f"COMPLETED! Processed {len(all_metrics)} datasets")
    print(f"Output directory: {session_dir}")
    print(f"{'='*80}")