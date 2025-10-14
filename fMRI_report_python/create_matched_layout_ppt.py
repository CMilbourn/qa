#!/usr/bin/env python3
"""
Analyze PowerPoint structure and recreate comprehensive report with same layout
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

def analyze_ppt_structure(ppt_path):
    """Analyze the structure of an existing PowerPoint presentation."""
    
    if not os.path.exists(ppt_path):
        print(f"❌ File not found: {ppt_path}")
        return None
    
    try:
        prs = Presentation(ppt_path)
        print(f"📊 Analyzing PowerPoint: {os.path.basename(ppt_path)}")
        print(f"Total slides: {len(prs.slides)}")
        print()
        
        structure = []
        
        for i, slide in enumerate(prs.slides):
            slide_info = {
                'slide_number': i + 1,
                'layout_name': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
                'layout_index': slide.slide_layout.slide_layout.layoutId if hasattr(slide.slide_layout, 'slide_layout') else -1,
                'shapes': []
            }
            
            # Analyze shapes on slide
            for shape in slide.shapes:
                shape_info = {
                    'type': str(shape.shape_type),
                    'name': shape.name if hasattr(shape, 'name') else 'Unknown'
                }
                
                # Get text content if it's a text shape
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text_content = shape.text_frame.text.strip()
                    if text_content:
                        shape_info['text'] = text_content[:100] + ('...' if len(text_content) > 100 else '')
                
                # Check if it's an image
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_info['is_image'] = True
                    if hasattr(shape, 'image'):
                        shape_info['image_format'] = shape.image.content_type if hasattr(shape.image, 'content_type') else 'Unknown'
                
                slide_info['shapes'].append(shape_info)
            
            structure.append(slide_info)
            
            # Print slide summary
            print(f"Slide {i+1}: {slide_info['layout_name']}")
            for shape_info in slide_info['shapes']:
                if 'text' in shape_info:
                    print(f"  - Text: {shape_info['text']}")
                elif shape_info.get('is_image'):
                    print(f"  - Image: {shape_info.get('image_format', 'Unknown format')}")
                else:
                    print(f"  - Shape: {shape_info['type']}")
            print()
        
        return structure
        
    except Exception as e:
        print(f"❌ Error analyzing PowerPoint: {e}")
        return None

def create_matching_layout_ppt(qa_parent_dir, output_file, reference_structure=None):
    """Create PowerPoint with layout matching the reference structure."""
    
    # Create presentation
    prs = Presentation()
    
    # Get all QA directories
    qa_dirs = []
    if os.path.exists(qa_parent_dir):
        for item in os.listdir(qa_parent_dir):
            item_path = os.path.join(qa_parent_dir, item)
            if os.path.isdir(item_path):
                qa_dirs.append((item, item_path))
    
    qa_dirs.sort()  # Sort alphabetically
    
    if not qa_dirs:
        print(f"❌ No QA directories found in {qa_parent_dir}")
        return False
    
    print(f"📊 Creating PowerPoint with {len(qa_dirs)} QA result sets")
    
    # Title slide (mimicking original structure)
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "fMRI Quality Assessment Report"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"sub001 Complete Analysis - {len(qa_dirs)} Acquisitions"
    
    # Summary slide with metrics table (like original)
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "QA Summary Overview"
    
    # Create summary text with metrics
    summary_text = "Acquisition Summary:\n\n"
    
    for i, (qa_name, qa_dir) in enumerate(qa_dirs, 1):
        summary_file = os.path.join(qa_dir, 'qa_summary.json')
        
        if os.path.exists(summary_file):
            with open(summary_file, 'r') as f:
                summary_info = json.load(f)
        else:
            summary_info = {}
        
        # Extract scan type from filename
        if "2mm" in qa_name and "1.5" not in qa_name:
            scan_type = "2mm Resolution"
        elif "2mm_1.5" in qa_name:
            scan_type = "2mm × 1.5mm"
        elif "1.5_mm_iso" in qa_name:
            scan_type = "1.5mm Isotropic"
        else:
            scan_type = "Unknown"
        
        tr = summary_info.get('tr', 'N/A')
        shape = summary_info.get('shape', [0,0,0,0])
        mean_tsnr = summary_info.get('mean_tsnr', 0)
        mean_isnr = summary_info.get('mean_isnr', 0)
        ernst_scaling = summary_info.get('ernst_scaling', 1.0)
        
        summary_text += f"{i}. {scan_type}\n"
        summary_text += f"   • Matrix: {shape[0]}×{shape[1]}×{shape[2]} × {shape[3]} timepoints\n"
        summary_text += f"   • TR: {tr}s (Ernst scaling: {ernst_scaling:.3f})\n"
        summary_text += f"   • tSNR: {mean_tsnr:.1f} | iSNR: {mean_isnr:.0f}\n\n"
    
    slide.placeholders[1].text = summary_text
    
    # Process each QA directory with consistent slide structure
    image_types = [
        ('mean_montage.png', 'Mean Signal Intensity'),
        ('tsnr_montage.png', 'Temporal SNR (tSNR)'),
        ('isnr_montage.png', 'Instantaneous SNR (iSNR)'),
        ('sample_timeseries.png', 'Time Series Analysis')
    ]
    
    for qa_name, qa_dir in qa_dirs:
        print(f"Processing: {qa_name}")
        
        # Load QA summary
        summary_file = os.path.join(qa_dir, 'qa_summary.json')
        if os.path.exists(summary_file):
            with open(summary_file, 'r') as f:
                summary_info = json.load(f)
        else:
            summary_info = {}
        
        # Section title slide for each acquisition
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        
        if "2mm" in qa_name and "1.5" not in qa_name:
            scan_title = "2mm Resolution Acquisition"
        elif "2mm_1.5" in qa_name:
            scan_title = "2mm × 1.5mm Resolution Acquisition"
        elif "1.5_mm_iso" in qa_name:
            scan_title = "1.5mm Isotropic Resolution Acquisition"
        else:
            scan_title = "fMRI Acquisition"
        
        slide.shapes.title.text = scan_title
        
        # Add acquisition parameters as subtitle
        tr = summary_info.get('tr', 'Unknown')
        shape = summary_info.get('shape', [0,0,0,0])
        ernst_scaling = summary_info.get('ernst_scaling', 1.0)
        mean_tsnr = summary_info.get('mean_tsnr', 0)
        mean_isnr = summary_info.get('mean_isnr', 0)
        
        subtitle = f"TR: {tr}s | Matrix: {shape[0]}×{shape[1]}×{shape[2]}×{shape[3]} | tSNR: {mean_tsnr:.1f} | iSNR: {mean_isnr:.0f}"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
        
        # Add image slides for each type (matching original layout)
        for image_file, image_title in image_types:
            image_path = os.path.join(qa_dir, image_file)
            
            if os.path.exists(image_path):
                # Use blank layout for image slides (like original)
                slide_layout = prs.slide_layouts[6]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title at top
                title_left = Inches(0.5)
                title_top = Inches(0.2)
                title_width = Inches(9)
                title_height = Inches(0.8)
                
                title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
                title_frame = title_shape.text_frame
                title_frame.text = f"{scan_title}: {image_title}"
                
                # Format title
                title_paragraph = title_frame.paragraphs[0]
                title_paragraph.font.size = Pt(24)
                title_paragraph.font.bold = True
                title_paragraph.alignment = PP_ALIGN.CENTER
                
                # Add image (preserve aspect ratio, centered)
                try:
                    img_left = Inches(1)
                    img_top = Inches(1.2)
                    img_width = Inches(8)
                    
                    slide.shapes.add_picture(image_path, img_left, img_top, width=img_width)
                    
                except Exception as e:
                    print(f"  ⚠️  Error adding image {image_file}: {e}")
                
                # Add metrics text at bottom for SNR images
                if 'snr' in image_file.lower():
                    metrics_left = Inches(0.5)
                    metrics_top = Inches(6.5)
                    metrics_width = Inches(9)
                    metrics_height = Inches(1)
                    
                    metrics_shape = slide.shapes.add_textbox(metrics_left, metrics_top, metrics_width, metrics_height)
                    metrics_frame = metrics_shape.text_frame
                    
                    if 'tsnr' in image_file.lower():
                        metrics_frame.text = f"Mean Temporal SNR: {mean_tsnr:.2f} (Ernst scaling: {ernst_scaling:.3f})"
                    elif 'isnr' in image_file.lower():
                        metrics_frame.text = f"Mean Instantaneous SNR: {mean_isnr:.0f} (Ernst scaling: {ernst_scaling:.3f})"
                    
                    # Format metrics text
                    metrics_paragraph = metrics_frame.paragraphs[0]
                    metrics_paragraph.font.size = Pt(14)
                    metrics_paragraph.alignment = PP_ALIGN.CENTER
                
                print(f"  ✅ Added: {image_title}")
            else:
                print(f"  ⚠️  Missing: {image_file}")
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✅ PowerPoint saved: {output_file}")
        
        # Get file size
        file_size = os.path.getsize(output_file) / (1024 * 1024)  # MB
        print(f"📊 File size: {file_size:.1f} MB")
        print(f"📄 Total slides: {len(prs.slides)}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error saving PowerPoint: {e}")
        return False

def main():
    """Main function."""
    
    # Analyze the original PowerPoint structure
    original_ppt = "/Users/cmilbourn/Documents/Sweet_Data/Development_Data/Sweet_Data_BIDS_Dev/sub001/sub001-visit001/func/sub001_QA_Report.pptx"
    
    print("🔍 Analyzing original PowerPoint structure...")
    structure = analyze_ppt_structure(original_ppt)
    
    print("\n" + "="*80)
    print("🚀 Creating new comprehensive report with matching layout...")
    print("="*80 + "\n")
    
    # Create new comprehensive report
    qa_parent_dir = "/Users/cmilbourn/Documents/GitHub/qa/fMRI_report_python/qa_output_batch"
    output_file = "/Users/cmilbourn/Documents/GitHub/qa/fMRI_report_python/sub001_comprehensive_QA_report_matched_layout.pptx"
    
    success = create_matching_layout_ppt(qa_parent_dir, output_file, structure)
    
    if success:
        print(f"🎉 Comprehensive report created with matching layout!")
        print(f"📁 Saved as: {os.path.basename(output_file)}")
    else:
        print(f"❌ Failed to create comprehensive report!")

if __name__ == "__main__":
    main()