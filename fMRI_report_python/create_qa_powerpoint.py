#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Create PowerPoint presentation from QA output images
"""

import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import sys

def create_qa_powerpoint(output_dir, subject_name="QA Analysis"):
    """
    Create a PowerPoint presentation from QA output images
    
    Parameters:
    -----------
    output_dir : str
        Path to the QA output directory containing PNG images
    subject_name : str
        Name to use in the title slide
    """
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define image order and titles
    image_config = [
        ('Mean_image.png', 'Mean Image'),
        ('mean_montage.png', 'Mean Image Montage - All Slices'),
        ('masked_noise.png', 'Noise Volume Analysis'),
        ('noise_volume_montage.png', 'Noise Volume Montage'),
        ('masked_noise_volume_montage.png', 'Masked Noise Volume Montage'),
        ('iSNR_sag.png', 'iSNR Map - Sagittal View'),
        ('iSNR_cor.png', 'iSNR Map - Coronal View'),
        ('isnr_montage.png', 'iSNR Montage - All Slices'),
        ('tSNR_sag.png', 'tSNR Map - Sagittal View'),
        ('tSNR_cor.png', 'tSNR Map - Coronal View'),
        ('tSNR_per_unit_time.png', 'tSNR per Unit Time'),
        ('tSNR_raw.png', 'Raw tSNR Map'),
        ('tSNR_montage.png', 'tSNR Montage - All Slices'),
        ('tSNR_w_ROI_images.png', 'tSNR with ROI'),
        ('TS_images.png', 'Time Series Analysis'),
        ('SSN.png', 'Static Spatial Noise'),
    ]
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "fMRI Quality Assurance Report"
    subtitle.text = f"{subject_name}\n{os.path.basename(output_dir)}"
    
    # Style title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Add slides for each image
    for img_filename, img_title in image_config:
        img_path = os.path.join(output_dir, img_filename)
        
        if not os.path.exists(img_path):
            print(f"Warning: Image not found: {img_filename}")
            continue
        
        # Create blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add title
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(9)
        height = Inches(0.6)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = img_title
        
        # Style title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        
        # Add image
        img_left = Inches(0.5)
        img_top = Inches(1.1)
        img_width = Inches(9)
        
        try:
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
        except Exception as e:
            print(f"Error adding image {img_filename}: {e}")
    
    # Save presentation
    pptx_filename = os.path.join(output_dir, f"QA_Report_{os.path.basename(output_dir)}.pptx")
    prs.save(pptx_filename)
    print(f"PowerPoint created: {pptx_filename}")
    
    return pptx_filename


def find_latest_qa_output(base_dir):
    """
    Find the most recently created QA output directory
    
    Parameters:
    -----------
    base_dir : str
        Base directory to search for qa_output_* folders
    
    Returns:
    --------
    str or None
        Path to the most recent qa_output directory
    """
    qa_dirs = glob.glob(os.path.join(base_dir, 'qa_output_*'))
    
    if not qa_dirs:
        return None
    
    # Sort by modification time, most recent first
    qa_dirs.sort(key=os.path.getmtime, reverse=True)
    
    return qa_dirs[0]


if __name__ == "__main__":
    if len(sys.argv) > 1:
        # Use provided directory
        output_dir = sys.argv[1]
        if not os.path.exists(output_dir):
            print(f"Error: Directory not found: {output_dir}")
            sys.exit(1)
    else:
        # Try to find the most recent QA output directory
        # Default search paths
        search_paths = [
            '/Users/cmilbourn/Documents/Sweet_Data/Development_Data/sub003/sub003-visit001-ses001/',
            os.getcwd(),
        ]
        
        output_dir = None
        for search_path in search_paths:
            if os.path.exists(search_path):
                output_dir = find_latest_qa_output(search_path)
                if output_dir:
                    break
        
        if not output_dir:
            print("Error: No QA output directory found.")
            print("Usage: python create_qa_powerpoint.py <qa_output_directory>")
            sys.exit(1)
    
    print(f"Creating PowerPoint from: {output_dir}")
    
    # Extract subject name from directory
    subject_name = os.path.basename(output_dir).replace('qa_output_', '').split('_2025')[0]
    
    create_qa_powerpoint(output_dir, subject_name)
